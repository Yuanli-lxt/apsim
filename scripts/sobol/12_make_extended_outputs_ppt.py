# -*- coding: utf-8 -*-
"""
Create a Chinese PPT report for the APSIM Classic Sobol extended-output analysis.

The deck uses native PowerPoint text boxes for all Chinese text and explicitly
sets Microsoft YaHei as both western and East Asian font.
"""

from __future__ import annotations

import json
import zipfile
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2] / "outputs" / "sobol" / "organized_outputs_screened_N128_extended_outputs_20260518_120007"
FINAL = ROOT / "final_results"
FIG = FINAL / "figures"
OUT_DIR = FINAL / "ppt"
PPTX = OUT_DIR / "APSIM_Sobol_extended_outputs_report_CN.pptx"
QA = OUT_DIR / "qa_report.md"
NOTES = OUT_DIR / "speaker_notes_cn.md"

FONT_CN = "Microsoft YaHei"
FONT_EN = "Aptos"

COLORS = {
    "ink": RGBColor(28, 42, 56),
    "muted": RGBColor(92, 107, 121),
    "light": RGBColor(244, 247, 250),
    "line": RGBColor(210, 218, 226),
    "blue": RGBColor(43, 91, 136),
    "teal": RGBColor(22, 128, 118),
    "orange": RGBColor(184, 111, 43),
    "red": RGBColor(166, 64, 64),
    "white": RGBColor(255, 255, 255),
}


def set_run_font(run, size=14, bold=False, color=None):
    run.font.name = FONT_CN
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    for tag in ["latin", "ea", "cs"]:
        child = rpr.find(qn(f"a:{tag}"))
        if child is None:
            child = OxmlElement(f"a:{tag}")
            rpr.append(child)
        child.set("typeface", FONT_CN)


def style_text_frame(tf, size=14, color=None, bold=False, align=None, line_spacing=1.08):
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for p in tf.paragraphs:
        if align is not None:
            p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            set_run_font(run, size=size, bold=bold, color=color)


def add_textbox(slide, x, y, w, h, text, size=14, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color or COLORS["ink"])
    if align is not None:
        p.alignment = align
    style_text_frame(tf, size=size, color=color or COLORS["ink"], bold=bold, align=align)
    return box


def add_bullets(slide, x, y, w, h, bullets, size=14, color=None, level0_bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        p.level = 0
        p.space_after = Pt(5)
        p.line_spacing = 1.08
        run = p.add_run()
        run.text = "• " + bullet
        set_run_font(run, size=size, bold=(level0_bold and i == 0), color=color or COLORS["ink"])
    style_text_frame(tf, size=size, color=color or COLORS["ink"])
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.48, 0.22, 12.2, 0.55, title, size=25, bold=True, color=COLORS["ink"])
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(0.88), Inches(12.3), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, 0.52, 0.92, 11.8, 0.26, subtitle, size=8.5, color=COLORS["muted"])


def add_footer(slide, source):
    add_textbox(slide, 0.55, 7.12, 12.2, 0.22, source, size=7.5, color=COLORS["muted"])


def add_panel(slide, x, y, w, h, fill=COLORS["light"], line=COLORS["line"]):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.7)
    return shape


def add_metric(slide, x, y, w, h, label, value, color=COLORS["blue"]):
    add_panel(slide, x, y, w, h, fill=COLORS["white"], line=COLORS["line"])
    add_textbox(slide, x + 0.10, y + 0.08, w - 0.20, 0.25, label, size=8.5, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.10, y + 0.35, w - 0.20, h - 0.38, value, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)


def add_image(slide, img_path: Path, x, y, w, h):
    if not img_path.exists():
        add_panel(slide, x, y, w, h, fill=RGBColor(255, 247, 247), line=COLORS["red"])
        add_textbox(slide, x + 0.1, y + h / 2 - 0.15, w - 0.2, 0.3, f"缺少图表：{img_path.name}", size=12, color=COLORS["red"], align=PP_ALIGN.CENTER)
        return None
    with Image.open(img_path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        new_w = w
        new_h = w / img_ratio
    else:
        new_h = h
        new_w = h * img_ratio
    left = x + (w - new_w) / 2
    top = y + (h - new_h) / 2
    return slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), width=Inches(new_w), height=Inches(new_h))


def add_table(slide, x, y, w, h, headers, rows, font_size=9.5):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["blue"]
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in cell.text_frame.paragraphs[0].runs:
            set_run_font(run, font_size, bold=True, color=COLORS["white"])
    for i, row in enumerate(rows, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["white"] if i % 2 else COLORS["light"]
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    set_run_font(run, font_size, color=COLORS["ink"])
    return table_shape


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def read_data():
    qc = pd.read_csv(FINAL / "extended_outputs_QC_summary.csv")
    top = pd.read_csv(FINAL / "ppt_ready_top10_ST_extended_outputs.csv")
    priority = pd.read_csv(FINAL / "ppt_ready_global_parameter_priority_extended_outputs.csv")
    with open(ROOT / "intermediate_and_raw_files" / "sobol_problem_definition.json", "r", encoding="utf-8") as f:
        problem = json.load(f)
    return qc, top, priority, problem


def top_params(top: pd.DataFrame, crop: str, target: str, n=4):
    g = top[(top["crop"] == crop) & (top["target"] == target)].head(n)
    return [(r["parameter"], f"{r['ST']:.3f}", f"{r['S1']:.3f}") for _, r in g.iterrows()]


def build_deck():
    qc, top, priority, problem = read_data()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    notes = []

    # Slide 1
    slide = blank_slide(prs)
    add_textbox(slide, 0.65, 0.55, 11.9, 0.95, "基于 APSIM Classic 的品种参数 Sobol 扩展输出分析", size=29, bold=True, color=COLORS["ink"])
    add_textbox(slide, 0.68, 1.55, 11.7, 0.35, "小麦与玉米：grain number、grain weight 与水分利用效率", size=15, color=COLORS["muted"])
    add_metric(slide, 0.75, 2.35, 2.25, 1.05, "APSIM", "Classic 7.10")
    add_metric(slide, 3.15, 2.35, 1.75, 1.05, "参数数 D", "13")
    add_metric(slide, 5.05, 2.35, 1.75, 1.05, "Sobol N", "128")
    add_metric(slide, 6.95, 2.35, 2.25, 1.05, "模拟数", "1920")
    add_metric(slide, 9.35, 2.35, 2.35, 1.05, "运行状态", "100% 完成", COLORS["teal"])
    add_panel(slide, 0.75, 4.05, 11.8, 1.65, fill=COLORS["light"])
    add_bullets(slide, 1.05, 4.28, 11.25, 1.05, [
        "本次 PPT 聚焦扩展输出变量，不替代主结果目录中的产量、生物量、LAI、开花期与成熟期分析。",
        "玉米产量构成变量改用大写 Grain 变量，解决小写变量输出为 ? 的问题。",
        "WUE 由 grain_yield 与实际 ET 计算，ET 采用 transpiration + soil evaporation。"
    ], size=14)
    add_footer(slide, f"结果目录：{ROOT}")
    notes.append("开场强调：这是主 N128 分析后的扩展输出模块，目的在于补齐产量构成和 WUE。")

    # Slide 2
    slide = blank_slide(prs)
    add_title(slide, "扩展分析先解决了 APSIM 输出变量可读性问题")
    add_bullets(slide, 0.65, 1.28, 4.0, 2.3, [
        "旧变量不可用：maize.grain_no、maize.grain_size、maize.grain_wt 输出为 ?。",
        "正式变量改为：maize.GrainNo 与 maize.GrainSize。",
        "maize.GrainSize 表示 1000 grain weight，数值约 341-527。",
        "变量来源可由 extended_variable_mapping_report.csv 追溯。"
    ], size=13.2)
    add_table(slide, 5.0, 1.18, 7.4, 2.55,
              ["标准变量", "APSIM 输出列", "用途"],
              [
                  ["maize grain_number", "MaizeGrainNoCapital", "籽粒数"],
                  ["maize grain_weight", "MaizeGrainSizeCapital", "千粒重"],
                  ["maize grain_yield", "MaizeYield / paddock.maize.yield", "产量"],
                  ["maize WUE", "Yield / (transpiration + es)", "水分利用效率"],
              ],
              font_size=8.8)
    add_panel(slide, 0.85, 4.35, 11.8, 1.45, fill=RGBColor(247, 250, 252))
    add_textbox(slide, 1.05, 4.55, 11.35, 0.55, "关键检查", size=16, bold=True, color=COLORS["blue"])
    add_bullets(slide, 1.05, 5.03, 11.2, 0.55, [
        "available_output_columns.csv 中不再出现旧小写 maize grain 列；大写 Grain 列进入输出并被汇总。"
    ], size=13)
    add_footer(slide, "依据文件：N128_extended_outputs_run_manifest.md；extended_variable_mapping_report.csv；available_output_columns.csv")
    notes.append("说明为什么先改 report/output：APSIM Classic 变量名具有模块和大小写差异，小写变量虽能生成列但没有有效值。")

    # Slide 3
    slide = blank_slide(prs)
    add_title(slide, "质量控制显示扩展变量可用于后续 Sobol 指数计算")
    add_metric(slide, 0.75, 1.25, 2.1, 1.0, "完成样本", "1920/1920", COLORS["teal"])
    add_metric(slide, 3.05, 1.25, 2.1, 1.0, "缺失值", "0", COLORS["teal"])
    add_metric(slide, 5.35, 1.25, 2.1, 1.0, "问号值 ?", "0", COLORS["teal"])
    add_metric(slide, 7.65, 1.25, 2.1, 1.0, "NaN", "0", COLORS["teal"])
    add_metric(slide, 9.95, 1.25, 2.1, 1.0, "常数列", "0", COLORS["teal"])
    qc_rows = []
    for crop, var in [
        ("maize", "grain_number"),
        ("maize", "grain_weight"),
        ("maize", "water_use_efficiency_yield"),
        ("wheat", "grain_number"),
        ("wheat", "grain_weight"),
        ("wheat", "water_use_efficiency_yield"),
    ]:
        row = qc[(qc["crop"] == crop) & (qc["variable"] == var)].iloc[0]
        qc_rows.append([crop, var, int(row["present_count"]), f"{row['min']:.3g}", f"{row['max']:.3g}"])
    add_table(slide, 0.75, 2.75, 11.85, 2.55, ["作物", "变量", "有效数", "最小值", "最大值"], qc_rows, font_size=9.2)
    add_bullets(slide, 0.95, 5.7, 11.5, 0.6, [
        "运行日志未检出 APSIM 批量运行错误；Wheat.xml 与 Maize.xml 运行后恢复检查无差异。"
    ], size=13)
    add_footer(slide, "依据文件：extended_outputs_QC_summary.csv/md；simulation_index.csv；05_run_apsim_batch.log")
    notes.append("这一页只讲数据质量，不讲生理机制。重点是扩展变量可以放心进入敏感性分析。")

    # Slide 4
    slide = blank_slide(prs)
    add_title(slide, "读图规则：ST 是本次汇报的主指标，S1 用作辅助")
    add_panel(slide, 0.8, 1.25, 3.6, 3.55, fill=COLORS["white"])
    add_textbox(slide, 1.05, 1.55, 3.1, 0.4, "S1：主效应", size=19, bold=True, color=COLORS["blue"], align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.08, 2.1, 3.0, 1.25, "单个参数独立造成的输出变化比例。S1 高说明该参数本身就能解释较多变化。", size=13, color=COLORS["ink"])
    add_panel(slide, 4.95, 1.25, 3.6, 3.55, fill=COLORS["white"])
    add_textbox(slide, 5.2, 1.55, 3.1, 0.4, "ST：总效应", size=19, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)
    add_textbox(slide, 5.23, 2.1, 3.0, 1.25, "参数自身影响加上与其他参数的交互影响。本汇报主要按 ST 排名。", size=13, color=COLORS["ink"])
    add_panel(slide, 9.1, 1.25, 3.3, 3.55, fill=COLORS["white"])
    add_textbox(slide, 9.35, 1.55, 2.8, 0.4, "ST - S1", size=19, bold=True, color=COLORS["orange"], align=PP_ALIGN.CENTER)
    add_textbox(slide, 9.38, 2.1, 2.75, 1.25, "差值越大，越可能存在交互或非线性路径。本轮未计算 S2，因此不定位具体两两交互。", size=13, color=COLORS["ink"])
    add_bullets(slide, 1.0, 5.35, 11.5, 0.65, [
        "N=128 对 D=13 的核心排序足够用于筛选；小效应参数、负 S1、S1>ST 不做强解释。"
    ], size=13.5)
    add_footer(slide, "依据文件：sobol_indices_summary.csv；sobol_problem_definition.json")
    notes.append("提醒听众：S1 是单独效应，ST 是总效应；ST 高但 S1 不高时，不要简单说主效应强。")

    # Slide 5
    slide = blank_slide(prs)
    add_title(slide, "总体热图显示：玉米偏热时间，小麦偏物候-春化-光周期")
    add_image(slide, FIG / "heatmap_parameter_by_output_ST.png", 0.55, 1.12, 8.45, 5.55)
    add_panel(slide, 9.35, 1.22, 3.3, 4.9, fill=COLORS["white"])
    top5 = priority.head(5)
    bullets = [f"{r.parameter_name}: 平均 ST {r.mean_ST:.3f}" for _, r in top5.iterrows()]
    add_textbox(slide, 9.62, 1.45, 2.85, 0.35, "全局平均 ST Top 5", size=15.5, bold=True, color=COLORS["blue"])
    add_bullets(slide, 9.55, 1.95, 2.9, 2.05, bullets, size=11.2)
    add_bullets(slide, 9.55, 4.35, 2.9, 1.05, [
        "玉米：tt_flag_to_flower 最突出。",
        "小麦：tt_floral_initiation、tt_end_of_juvenile、photop_sens、vern_sens 更集中。"
    ], size=11.2)
    add_footer(slide, "图：figures\\heatmap_parameter_by_output_ST.png；表：ppt_ready_global_parameter_priority_extended_outputs.csv")
    notes.append("热图用于建立整体印象：不同作物的敏感性结构不同，不要只看某一个输出变量。")

    # Slide 6
    slide = blank_slide(prs)
    add_title(slide, "玉米产量构成对 tt_flag_to_flower 高度敏感")
    add_image(slide, FIG / "ST_grain_number_maize_P01_shandong_2025_v503_joint_iter60.png", 0.45, 1.1, 7.5, 5.55)
    maize_gn = top_params(top, "maize", "grain_number", 4)
    maize_gw = top_params(top, "maize", "grain_weight", 4)
    add_panel(slide, 8.25, 1.15, 4.5, 2.25, fill=COLORS["white"])
    add_textbox(slide, 8.48, 1.35, 4.05, 0.3, "grain_number：ST 前四", size=14.5, bold=True, color=COLORS["blue"])
    add_table(slide, 8.45, 1.75, 3.95, 1.25, ["参数", "ST", "S1"], maize_gn, font_size=7.8)
    add_panel(slide, 8.25, 3.75, 4.5, 2.25, fill=COLORS["white"])
    add_textbox(slide, 8.48, 3.95, 4.05, 0.3, "grain_weight：ST 前四", size=14.5, bold=True, color=COLORS["blue"])
    add_table(slide, 8.45, 4.35, 3.95, 1.25, ["参数", "ST", "S1"], maize_gw, font_size=7.8)
    add_footer(slide, "图：ST_grain_number_maize_*.png；表：ppt_ready_top10_ST_extended_outputs.csv")
    notes.append("玉米 grain_number 与 grain_weight 都由 tt_flag_to_flower 主导；grain_weight 中 tt_flower_to_maturity 也有明显贡献。")

    # Slide 7
    slide = blank_slide(prs)
    add_title(slide, "玉米 WUE 中存在更明显的交互/非线性信号")
    add_image(slide, FIG / "S1_vs_ST_water_use_efficiency_biomass_maize_P01_shandong_2025_v503_joint_iter60.png", 0.55, 1.13, 7.6, 5.45)
    maize_wue = top_params(top, "maize", "water_use_efficiency_biomass", 5)
    add_panel(slide, 8.55, 1.25, 3.9, 4.65, fill=COLORS["white"])
    add_textbox(slide, 8.8, 1.48, 3.4, 0.3, "biomass WUE：ST 前五", size=14.5, bold=True, color=COLORS["teal"])
    add_table(slide, 8.75, 1.88, 3.35, 1.55, ["参数", "ST", "S1"], maize_wue, font_size=7.5)
    add_bullets(slide, 8.72, 3.8, 3.45, 1.3, [
        "tt_flag_to_flower：ST 0.555，S1 0.123，交互贡献大。",
        "tt_endjuv_to_init 与 tt_flower_to_maturity 同样重要。",
        "WUE_biomass 比 WUE_yield 更受多参数共同影响。"
    ], size=10.8)
    add_footer(slide, "图：S1_vs_ST_water_use_efficiency_biomass_maize_*.png；表：sobol_indices_summary.csv")
    notes.append("这里突出 ST 与 S1 的差异：不是单参数主效应，而是多物候参数共同改变生物量与耗水。")

    # Slide 8
    slide = blank_slide(prs)
    add_title(slide, "小麦籽粒数主要受花芽分化、幼年期结束和光周期控制")
    add_image(slide, FIG / "ST_grain_number_wheat_Jimai70_v132_joint_iter353.png", 0.45, 1.1, 7.55, 5.55)
    wheat_gn = top_params(top, "wheat", "grain_number", 5)
    wheat_gw = top_params(top, "wheat", "grain_weight", 5)
    add_panel(slide, 8.25, 1.12, 4.5, 2.3, fill=COLORS["white"])
    add_textbox(slide, 8.48, 1.34, 4.05, 0.3, "grain_number：ST 前五", size=14.5, bold=True, color=COLORS["blue"])
    add_table(slide, 8.45, 1.74, 3.95, 1.4, ["参数", "ST", "S1"], wheat_gn, font_size=7.6)
    add_panel(slide, 8.25, 3.72, 4.5, 2.3, fill=COLORS["white"])
    add_textbox(slide, 8.48, 3.94, 4.05, 0.3, "grain_weight：ST 前五", size=14.5, bold=True, color=COLORS["blue"])
    add_table(slide, 8.45, 4.34, 3.95, 1.4, ["参数", "ST", "S1"], wheat_gw, font_size=7.6)
    add_footer(slide, "图：ST_grain_number_wheat_*.png；表：ppt_ready_top10_ST_extended_outputs.csv")
    notes.append("小麦与玉米不同：没有一个参数压倒所有变量，而是多种物候响应共同控制。")

    # Slide 9
    slide = blank_slide(prs)
    add_title(slide, "小麦 WUE 对物候与籽粒大小均有响应，但不确定性更高")
    add_image(slide, FIG / "ST_water_use_efficiency_biomass_wheat_Jimai70_v132_joint_iter353.png", 0.55, 1.12, 7.5, 5.55)
    wheat_wue = top_params(top, "wheat", "water_use_efficiency_biomass", 6)
    add_panel(slide, 8.35, 1.15, 4.25, 2.55, fill=COLORS["white"])
    add_textbox(slide, 8.58, 1.38, 3.8, 0.3, "biomass WUE：ST 前六", size=14.5, bold=True, color=COLORS["teal"])
    add_table(slide, 8.55, 1.78, 3.65, 1.55, ["参数", "ST", "S1"], wheat_wue, font_size=7.3)
    add_panel(slide, 8.35, 4.05, 4.25, 1.65, fill=RGBColor(255, 250, 242), line=RGBColor(230, 190, 130))
    add_textbox(slide, 8.58, 4.25, 3.8, 0.3, "限制提醒", size=14, bold=True, color=COLORS["orange"])
    add_bullets(slide, 8.55, 4.62, 3.75, 0.85, [
        "多个参数 ST_conf/ST > 0.5。",
        "结果适合说明趋势，不宜过度精确排序。"
    ], size=10.6, color=COLORS["ink"])
    add_footer(slide, "图：ST_water_use_efficiency_biomass_wheat_*.png；质量标记来自 sobol_indices_summary.csv")
    notes.append("小麦 WUE_biomass 是最需要谨慎的一页：可讲参数组重要性，少讲精确名次。")

    # Slide 10
    slide = blank_slide(prs)
    add_title(slide, "扩展输出分析给出校准优先级：先调物候，再看作物特异参数")
    add_image(slide, FIG / "crop_comparison_mean_ST.png", 0.55, 1.05, 6.95, 5.7)
    add_panel(slide, 7.85, 1.1, 4.85, 4.95, fill=COLORS["white"])
    add_textbox(slide, 8.1, 1.35, 4.4, 0.32, "可直接带走的三点", size=16, bold=True, color=COLORS["blue"])
    add_bullets(slide, 8.08, 1.85, 4.35, 2.1, [
        "玉米：优先校准 tt_flag_to_flower，其次 tt_endjuv_to_init 与 tt_flower_to_maturity。",
        "小麦：优先关注 tt_floral_initiation、tt_end_of_juvenile、photop_sens、vern_sens。",
        "低优先级参数：largestLeafParams[1/2] 与 rue 在本扩展变量中贡献较低。"
    ], size=12.1)
    add_panel(slide, 8.05, 4.55, 4.35, 0.85, fill=COLORS["light"])
    add_textbox(slide, 8.25, 4.72, 3.95, 0.38, "后续：与主 N128 结果合并，形成最终校准参数清单。", size=12.5, bold=True, color=COLORS["teal"])
    add_footer(slide, "图：crop_comparison_mean_ST.png；表：ppt_ready_global_parameter_priority_extended_outputs.csv")
    notes.append("结束时强调：本结果不是替代主结果，而是补充产量构成和 WUE 的校准依据。")

    # Apply background and remove accidental default fills.
    for slide in prs.slides:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(252, 253, 254)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX)

    NOTES.write_text(
        "\n\n".join([f"## 第 {i+1} 页\n{note}" for i, note in enumerate(notes)]),
        encoding="utf-8",
    )

    # Verify by reopening.
    reopened = Presentation(str(PPTX))
    chinese_count = 0
    question_count = 0
    with zipfile.ZipFile(PPTX) as zf:
        for name in zf.namelist():
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                txt = zf.read(name).decode("utf-8", errors="ignore")
                chinese_count += sum("\u4e00" <= ch <= "\u9fff" for ch in txt)
                question_count += txt.count("????")
    qa = [
        "# PPT QA Report",
        f"PPTX: {PPTX}",
        f"slide_count: {len(reopened.slides)}",
        f"embedded_media_count: {len(list((OUT_DIR).glob('*')))}",
        f"chinese_characters_in_slide_xml: {chinese_count}",
        f"question_mark_sequences_in_slide_xml: {question_count}",
        "font: Microsoft YaHei set for native text runs, including East Asian font.",
        "figures: sourced from final_results/figures PNG files.",
        "notes: speaker_notes_cn.md",
    ]
    QA.write_text("\n".join(qa), encoding="utf-8")
    if chinese_count <= 0 or question_count > 0:
        raise RuntimeError("PPTX Chinese verification failed.")


if __name__ == "__main__":
    build_deck()
    print(PPTX)
