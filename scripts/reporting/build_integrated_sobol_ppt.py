# -*- coding: utf-8 -*-
"""Build an integrated Chinese PPT report for cultivar and system Sobol results."""

from __future__ import annotations

import json
import re
import shutil
import zipfile
from datetime import datetime
from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from matplotlib.font_manager import FontProperties
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[2]
SYSTEM_DIR = ROOT / "outputs/system_sensitivity/final_results"
CULTIVAR_DIR = ROOT / "outputs/sobol/organized_outputs_screened_N128_20260515_185604/final_results"
BEST_DIR = ROOT / "results/local_sobol_guided_search_20260519_output_sobol/best"

OUT_DIR = SYSTEM_DIR / f"ppt_report_integrated_sobol_cn_fixed_{datetime.now():%Y%m%d_%H%M%S}"
ASSET_DIR = OUT_DIR / "assets"

FONT_CANDIDATES = [
    Path(r"C:\Windows\Fonts\msyh.ttc"),
    Path(r"C:\Windows\Fonts\Noto Sans SC (TrueType).otf"),
    Path(r"C:\Windows\Fonts\SourceHanSansCN-Normal.ttf"),
    Path(r"C:\Windows\Fonts\simhei.ttf"),
    Path(r"C:\Windows\Fonts\simsun.ttc"),
]
FONT_PATH = next((p for p in FONT_CANDIDATES if p.exists()), None)
FONT_NAME = "Microsoft YaHei"

BG = RGBColor(248, 249, 251)
INK = RGBColor(31, 41, 55)
MUTED = RGBColor(92, 104, 121)
BLUE = RGBColor(42, 93, 168)
GREEN = RGBColor(43, 128, 94)
ORANGE = RGBColor(199, 111, 38)
LIGHT = RGBColor(232, 238, 247)


def ensure_inputs() -> None:
    required = [
        SYSTEM_DIR / "sobol_indices_summary.csv",
        SYSTEM_DIR / "sobol_N64_vs_N128_stability_report.csv",
        CULTIVAR_DIR / "sobol_indices_summary.csv",
        BEST_DIR / "best_selection.json",
        BEST_DIR / "metrics.json",
    ]
    missing = [str(p) for p in required if not p.exists()]
    if missing:
        raise FileNotFoundError("缺少输入文件：\n" + "\n".join(missing))


def east_asia_font(run, size=None, bold=False, color=INK) -> None:
    run.font.name = FONT_NAME
    run.font.bold = bold
    run.font.color.rgb = color
    if size is not None:
        run.font.size = Pt(size)
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = r_pr.find(tag, namespaces=r_pr.nsmap)
        if node is None:
            node = OxmlElement(tag)
            r_pr.append(node)
        node.set("typeface", FONT_NAME)


def set_text_frame(tf, paragraphs) -> None:
    tf.clear()
    for i, item in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        p.level = item.get("level", 0)
        p.space_after = Pt(5)
        p.line_spacing = 1.12
        r = p.add_run()
        r.text = item["text"]
        east_asia_font(r, item.get("size", 16), item.get("bold", False), item.get("color", INK))


def add_textbox(slide, left, top, width, height, paragraphs, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or RGBColor(255, 255, 255)
    shape.line.color.rgb = line or RGBColor(220, 226, 235)
    shape.text_frame.margin_left = Cm(0.35)
    shape.text_frame.margin_right = Cm(0.35)
    shape.text_frame.margin_top = Cm(0.24)
    shape.text_frame.margin_bottom = Cm(0.18)
    set_text_frame(shape.text_frame, paragraphs)
    return shape


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Cm(0.85), Cm(0.45), Cm(31.8), Cm(1.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    east_asia_font(r, 26, True, INK)
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        east_asia_font(r2, 12, False, MUTED)
    return box


def add_footer(slide, source):
    box = slide.shapes.add_textbox(Cm(0.9), Cm(18.38), Cm(31.2), Cm(0.38))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = source
    east_asia_font(r, 8.5, False, MUTED)


def add_stat(slide, x, y, number, label, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(5.1), Cm(2.0))
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(220, 226, 235)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(number)
    east_asia_font(r, 25, True, color)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    east_asia_font(r2, 10.5, False, MUTED)


def param_label(key: str) -> str:
    mapping = {
        "soilwater__CN2Bare__scalar": "裸地径流曲线数 CN2Bare",
        "soilwater__CNRed__scalar": "覆盖减径流系数 CNRed",
        "soilwater__CNCov__scalar": "覆盖阈值 CNCov",
        "soilwater__DiffusConst__scalar": "扩散常数 DiffusConst",
        "soilwater__DiffusSlope__scalar": "扩散斜率 DiffusSlope",
        "soilwater__SWCON__all_layers_multiplier": "排水系数 SWCON",
        "surfaceom__RootCN__scalar": "根残体碳氮比 RootCN",
        "surfaceom__SoilCN__scalar": "土壤碳氮比 SoilCN",
        "soiln__FBiom__scalar": "微生物库比例 FBiom",
        "soilcrop__wheat__KL__all_layers_multiplier": "小麦根系取水 KL",
        "soilcrop__wheat__XF__all_layers_multiplier": "小麦根系探索 XF",
        "soilcrop__maize__XF__all_layers_multiplier": "玉米根系探索 XF",
    }
    if key in mapping:
        return mapping[key]
    crop = "小麦" if key.startswith("wheat__") else "玉米" if key.startswith("maize__") else ""
    parts = key.split("__")
    if len(parts) >= 3:
        return f"{crop} {parts[2]}".strip()
    return key


def prepare_data():
    sys_df = pd.read_csv(SYSTEM_DIR / "sobol_indices_summary.csv")
    cul_df = pd.read_csv(CULTIVAR_DIR / "sobol_indices_summary.csv")
    stab_df = pd.read_csv(SYSTEM_DIR / "sobol_N64_vs_N128_stability_report.csv")
    with open(BEST_DIR / "best_selection.json", "r", encoding="utf-8") as f:
        best = json.load(f)
    with open(BEST_DIR / "metrics.json", "r", encoding="utf-8") as f:
        metrics = json.load(f)
    return sys_df, cul_df, stab_df, best, metrics


def plot_bar(df, value_col, label_col, title, out_path, color="#2A5DA8", top=10):
    data = df.sort_values(value_col, ascending=False).head(top).copy()
    data = data.sort_values(value_col)
    font = FontProperties(fname=str(FONT_PATH)) if FONT_PATH else None
    plt.figure(figsize=(8.2, 4.6), dpi=180)
    plt.barh(data[label_col], data[value_col], color=color)
    plt.xlabel(value_col, fontproperties=font)
    plt.title(title, fontproperties=font, fontsize=13)
    plt.xticks(fontproperties=font, fontsize=8)
    plt.yticks(fontproperties=font, fontsize=8)
    plt.grid(axis="x", alpha=0.22)
    plt.tight_layout()
    plt.savefig(out_path, bbox_inches="tight")
    plt.close()


def build_figures(sys_df, cul_df, stab_df):
    ASSET_DIR.mkdir(parents=True, exist_ok=True)

    sys_top = (
        sys_df.groupby("parameter_key", as_index=False)["ST"]
        .mean()
        .assign(label=lambda d: d["parameter_key"].map(param_label))
    )
    plot_bar(sys_top, "ST", "label", "系统 Sobol：平均总效应指数 Top 参数", ASSET_DIR / "system_top_st.png", "#2A5DA8", 10)

    cul_top = (
        cul_df.groupby("parameter_key", as_index=False)["ST"]
        .mean()
        .assign(label=lambda d: d["parameter_key"].map(param_label))
    )
    plot_bar(cul_top, "ST", "label", "品种 Sobol：平均总效应指数 Top 参数", ASSET_DIR / "cultivar_top_st.png", "#2B805E", 10)

    crop_sys = (
        sys_df.groupby(["crop", "parameter_key"], as_index=False)["ST"]
        .mean()
        .sort_values(["crop", "ST"], ascending=[True, False])
    )
    crop_sys["label"] = crop_sys["parameter_key"].map(param_label)
    selected = crop_sys.groupby("crop").head(6)
    font = FontProperties(fname=str(FONT_PATH)) if FONT_PATH else None
    plt.figure(figsize=(8.2, 4.8), dpi=180)
    for i, crop in enumerate(["wheat", "maize"]):
        sub = selected[selected["crop"] == crop].sort_values("ST")
        plt.subplot(1, 2, i + 1)
        plt.barh(sub["label"], sub["ST"], color="#2A5DA8" if crop == "wheat" else "#C76F26")
        plt.title("小麦" if crop == "wheat" else "玉米", fontproperties=font)
        plt.xlabel("ST", fontproperties=font)
        plt.xticks(fontproperties=font, fontsize=7)
        plt.yticks(fontproperties=font, fontsize=7)
        plt.grid(axis="x", alpha=0.2)
    plt.suptitle("系统 Sobol：作物分组 Top 参数", fontproperties=font, fontsize=13)
    plt.tight_layout()
    plt.savefig(ASSET_DIR / "system_crop_split.png", bbox_inches="tight")
    plt.close()

    stable = (
        stab_df.groupby(["crop", "target_variable"], as_index=False)
        .agg(top5_overlap=("top5_overlap_count", "max"), spearman=("spearman_ST_rank", "max"), mean_abs_st=("mean_abs_ST_diff", "max"))
    )
    stable["label"] = stable["crop"] + " / " + stable["target_variable"]
    stable = stable.sort_values(["top5_overlap", "spearman"], ascending=False).head(12)
    font = FontProperties(fname=str(FONT_PATH)) if FONT_PATH else None
    plt.figure(figsize=(8.2, 4.4), dpi=180)
    plt.scatter(stable["spearman"], stable["top5_overlap"], s=90, color="#2B805E")
    for _, row in stable.iterrows():
        plt.text(row["spearman"] + 0.003, row["top5_overlap"] + 0.02, row["label"], fontsize=7, fontproperties=font)
    plt.xlabel("N64 与 N128 的 ST 排名 Spearman 相关", fontproperties=font)
    plt.ylabel("Top5 重叠个数", fontproperties=font)
    plt.title("Sobol 稳定性：N64 到 N128", fontproperties=font, fontsize=13)
    plt.xlim(max(0, stable["spearman"].min() - 0.05), 1.02)
    plt.ylim(0, 5.4)
    plt.grid(alpha=0.25)
    plt.tight_layout()
    plt.savefig(ASSET_DIR / "stability_n64_n128.png", bbox_inches="tight")
    plt.close()

    heatmap_src = SYSTEM_DIR / "figures/heatmap_parameter_by_output_ST.png"
    heatmap_dst = ASSET_DIR / "system_heatmap_parameter_by_output_ST.png"
    if heatmap_src.exists():
        shutil.copy2(heatmap_src, heatmap_dst)

    return {
        "system_top": ASSET_DIR / "system_top_st.png",
        "cultivar_top": ASSET_DIR / "cultivar_top_st.png",
        "crop_split": ASSET_DIR / "system_crop_split.png",
        "stability": ASSET_DIR / "stability_n64_n128.png",
        "heatmap": heatmap_dst if heatmap_dst.exists() else None,
        "sys_top_df": sys_top.sort_values("ST", ascending=False),
        "cul_top_df": cul_top.sort_values("ST", ascending=False),
    }


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def add_picture(slide, path, left, top, width=None, height=None):
    if path and Path(path).exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def build_ppt(sys_df, cul_df, stab_df, best, metrics, figs):
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    slide = blank_slide(prs)
    add_title(slide, "APSIM 模型 Sobol 敏感性分析汇总", "从品种参数校准到系统过程参数识别")
    add_textbox(slide, Cm(1.1), Cm(3.0), Cm(17.2), Cm(4.3), [
        {"text": "汇报对象：品种 Sobol、局部迭代最佳基准、系统 Sobol N128", "size": 18, "bold": True},
        {"text": "核心问题：在品种参数已经优化之后，哪些土壤水分、根系取水和氮素过程参数仍然主导模型输出。", "size": 15},
    ], fill=RGBColor(255, 255, 255))
    add_stat(slide, 20.0, 3.0, "N=128", "系统 Sobol 样本")
    add_stat(slide, 25.6, 3.0, "12", "系统参数组")
    add_stat(slide, 20.0, 5.35, "品种 + 系统", "两阶段整合")
    add_stat(slide, 25.6, 5.35, "N64/N128", "稳定性检查")
    add_footer(slide, "数据来源：cultivar Sobol N128、system Sobol N128、N64/N128 stability report")

    slide = blank_slide(prs)
    add_title(slide, "分析路线", "先固定品种最优基准，再分析系统过程参数")
    x0, y0 = 1.3, 3.0
    steps = [
        ("1 品种 Sobol", "识别物候、光周期、春化、RUE 等品种参数影响"),
        ("2 局部迭代", "将最佳品种组合固化为新的 calibrated baseline"),
        ("3 Morris 筛选", "从系统候选参数中筛出 12 个关键参数组"),
        ("4 Sobol 精算", "N64 试算，N128 稳定性检查，输出 S1 与 ST"),
    ]
    for i, (head, body) in enumerate(steps):
        add_textbox(slide, Cm(x0 + i * 8.0), Cm(y0), Cm(7.1), Cm(4.7), [
            {"text": head, "size": 18, "bold": True, "color": BLUE},
            {"text": body, "size": 13},
        ], fill=RGBColor(255, 255, 255))
    add_textbox(slide, Cm(2.0), Cm(10.2), Cm(29.5), Cm(3.2), [
        {"text": "固定项", "size": 17, "bold": True, "color": GREEN},
        {"text": "真实气象、真实初始水分、真实初始 NO3/NH4、非研究目标管理记录，以及已经校准的品种参数。", "size": 15},
        {"text": "因此下一阶段的结论应表述为：围绕校准后 APSIM 模型的系统敏感性分析。", "size": 15},
    ], fill=RGBColor(255, 255, 255))

    slide = blank_slide(prs)
    add_title(slide, "品种 Sobol 结果", "主要敏感性集中在物候时长、春化/光周期和 RUE")
    add_picture(slide, figs["cultivar_top"], Cm(1.0), Cm(2.2), width=Cm(18.0))
    add_textbox(slide, Cm(20.2), Cm(2.5), Cm(11.8), Cm(9.6), [
        {"text": "解读", "size": 18, "bold": True, "color": GREEN},
        {"text": "玉米侧重 tt_endjuv_to_init、tt_flower_to_maturity、tt_flag_to_flower 和 RUE。", "size": 14},
        {"text": "小麦侧重 tt_end_of_juvenile、photop_sens、vern_sens 和灌浆起始热时间。", "size": 14},
        {"text": "这些参数决定作物发育时钟和资源利用框架，因此已被固化到新基准模型。", "size": 14},
    ], fill=RGBColor(255, 255, 255))
    add_footer(slide, "Source: outputs/sobol/organized_outputs_screened_N128_20260515_185604/final_results/sobol_indices_summary.csv")

    slide = blank_slide(prs)
    add_title(slide, "校准后基准模型", "系统分析使用品种 Sobol 与局部迭代后的最佳模型")
    best_case = best.get("best_case", best)
    items = [
        f"最佳阶段：{best_case.get('phase', best.get('phase', 'C_fraction_after_maize'))}",
        f"case_id：{best_case.get('case_id', best.get('case_id', '20'))}",
        f"FractionFull：{best_case.get('FractionFull', best.get('FractionFull', 0.548))}",
        f"crit_fr_asw：{best_case.get('crit_fr_asw', best.get('crit_fr_asw', 0.56))}",
        "玉米 tt_flag_to_flower：38.3 调整到 38.683",
    ]
    add_textbox(slide, Cm(1.2), Cm(2.5), Cm(14.8), Cm(10.8), [{"text": "最佳基准信息", "size": 18, "bold": True, "color": BLUE}] + [{"text": x, "size": 15} for x in items], fill=RGBColor(255, 255, 255))
    add_textbox(slide, Cm(17.2), Cm(2.5), Cm(14.8), Cm(10.8), [
        {"text": "为什么这样做", "size": 18, "bold": True, "color": GREEN},
        {"text": "系统 Sobol 的目的不是重新搜索品种参数，而是在已经校准好的品种背景下，识别水分、根系和氮素过程的不确定性。", "size": 15},
        {"text": "这样可以避免品种物候误差掩盖土壤过程参数的影响。", "size": 15},
        {"text": "最终结论应避免写成完全外推结论，而应写成围绕校准模型的敏感性结论。", "size": 15},
    ], fill=RGBColor(255, 255, 255))
    add_footer(slide, "Source: results/local_sobol_guided_search_20260519_output_sobol/best/*.json")

    slide = blank_slide(prs)
    add_title(slide, "系统 Sobol 参数范围", "Morris 筛选后纳入 12 个系统参数组")
    params = [
        "CN2Bare、CNRed、CNCov",
        "DiffusConst、DiffusSlope、SWCON",
        "RootCN、SoilCN、FBiom",
        "wheat KL、wheat XF、maize XF",
    ]
    add_textbox(slide, Cm(1.3), Cm(2.4), Cm(13.8), Cm(7.5), [{"text": "参数组", "size": 18, "bold": True, "color": BLUE}] + [{"text": p, "size": 16} for p in params], fill=RGBColor(255, 255, 255))
    add_textbox(slide, Cm(16.3), Cm(2.4), Cm(15.8), Cm(7.5), [
        {"text": "扰动逻辑", "size": 18, "bold": True, "color": GREEN},
        {"text": "土壤水分过程：径流、扩散、排水。", "size": 15},
        {"text": "根系取水：按作物和土层整体乘数扰动 KL 或 XF。", "size": 15},
        {"text": "氮素/有机质：作为二级候选保留，用于检验系统输出是否敏感。", "size": 15},
    ], fill=RGBColor(255, 255, 255))
    add_stat(slide, 4.0, 11.0, "17", "Morris 候选")
    add_stat(slide, 10.0, 11.0, "12", "Sobol 参数")
    add_stat(slide, 16.0, 11.0, "1792", "N128 模拟")
    add_stat(slide, 22.0, 11.0, "S1/ST", "主要指数")

    slide = blank_slide(prs)
    add_title(slide, "系统 Sobol 主结果", "小麦 KL 与 SoilWater 扩散/排水参数是主导因子")
    add_picture(slide, figs["system_top"], Cm(1.0), Cm(2.2), width=Cm(18.3))
    add_textbox(slide, Cm(20.3), Cm(2.4), Cm(11.7), Cm(10.2), [
        {"text": "关键结论", "size": 18, "bold": True, "color": BLUE},
        {"text": "小麦 KL 的平均 ST 最高，说明根系取水能力是系统输出的核心不确定来源。", "size": 14},
        {"text": "DiffusSlope、DiffusConst、SWCON 共同控制水分再分配和排水过程。", "size": 14},
        {"text": "CN2Bare 影响径流分配，对产量、水分效率和水分收支有次级贡献。", "size": 14},
    ], fill=RGBColor(255, 255, 255))
    add_footer(slide, "Source: outputs/system_sensitivity/final_results/sobol_indices_summary.csv")

    slide = blank_slide(prs)
    add_title(slide, "作物分组差异", "系统敏感性更强地体现为小麦季水分限制")
    add_picture(slide, figs["crop_split"], Cm(1.0), Cm(2.2), width=Cm(19.0))
    add_textbox(slide, Cm(21.0), Cm(2.7), Cm(10.8), Cm(8.8), [
        {"text": "作物差异", "size": 18, "bold": True, "color": ORANGE},
        {"text": "小麦输出对 KL 的响应最强，说明冬小麦季的根系取水设定会明显改变最终产量和生物量。", "size": 14},
        {"text": "玉米输出更多受 SoilWater 扩散、排水和径流参数影响。", "size": 14},
        {"text": "这也支持后续把小麦 KL 作为重点校准或独立验证对象。", "size": 14},
    ], fill=RGBColor(255, 255, 255))

    slide = blank_slide(prs)
    add_title(slide, "输出指标响应", "grain、WUE、水分收支与物候指标需要分层解读")
    if figs["heatmap"]:
        add_picture(slide, figs["heatmap"], Cm(1.1), Cm(2.0), width=Cm(18.4))
    add_textbox(slide, Cm(20.1), Cm(2.2), Cm(12.0), Cm(10.8), [
        {"text": "指标说明", "size": 18, "bold": True, "color": GREEN},
        {"text": "grain_yield、biomass、LAI、grain_number、grain_weight 与 WUE 是主要解释对象。", "size": 14},
        {"text": "WUE 推荐按产量或生物量除以蒸散量计算；蒸散量可由 transpiration 加 soil_evaporation 得到。", "size": 14},
        {"text": "flowering_date 与 maturity_date 已加入输出；在系统参数扰动下若方差很小，Sobol 指数会接近零，这是合理结果。", "size": 14},
    ], fill=RGBColor(255, 255, 255))

    slide = blank_slide(prs)
    add_title(slide, "N64 到 N128 稳定性", "Top 参数总体稳定，N128 可作为正式结果")
    add_picture(slide, figs["stability"], Cm(1.1), Cm(2.2), width=Cm(18.8))
    overlap = stab_df["top5_overlap_count"].max()
    spearman = stab_df["spearman_ST_rank"].max()
    add_textbox(slide, Cm(21.0), Cm(2.5), Cm(10.7), Cm(8.6), [
        {"text": "稳定性判断", "size": 18, "bold": True, "color": BLUE},
        {"text": f"Top5 最大重叠数：{overlap:g} / 5", "size": 15},
        {"text": f"ST 排名 Spearman 最高值：{spearman:.3f}", "size": 15},
        {"text": "多数主要输出的 Top 参数在 N64 与 N128 之间保持一致，说明 N128 结果可用于正式整理。", "size": 15},
    ], fill=RGBColor(255, 255, 255))
    add_footer(slide, "Source: outputs/system_sensitivity/final_results/sobol_N64_vs_N128_stability_report.csv")

    slide = blank_slide(prs)
    add_title(slide, "两阶段结果的整合解释", "品种参数定时钟，系统参数定水分强度")
    add_textbox(slide, Cm(1.2), Cm(2.4), Cm(9.6), Cm(9.4), [
        {"text": "品种 Sobol", "size": 18, "bold": True, "color": GREEN},
        {"text": "控制发育阶段、光周期/春化响应和 RUE。", "size": 15},
        {"text": "决定模型是否在正确时间进入开花、灌浆和成熟阶段。", "size": 15},
    ], fill=RGBColor(255, 255, 255))
    add_textbox(slide, Cm(12.1), Cm(2.4), Cm(9.6), Cm(9.4), [
        {"text": "系统 Sobol", "size": 18, "bold": True, "color": BLUE},
        {"text": "控制根系取水、水分扩散、排水和径流。", "size": 15},
        {"text": "决定校准后模型在不同水分条件下的产量和水分效率响应。", "size": 15},
    ], fill=RGBColor(255, 255, 255))
    add_textbox(slide, Cm(23.0), Cm(2.4), Cm(9.0), Cm(9.4), [
        {"text": "综合结论", "size": 18, "bold": True, "color": ORANGE},
        {"text": "先校准品种参数，再分析系统参数，是这套 APSIM 敏感性分析更稳妥的顺序。", "size": 15},
        {"text": "论文中可强调这是围绕 calibrated baseline 的局部全局敏感性分析。", "size": 15},
    ], fill=RGBColor(255, 255, 255))

    slide = blank_slide(prs)
    add_title(slide, "下一步建议", "从敏感性排序走向校准、验证和论文表达")
    add_textbox(slide, Cm(1.2), Cm(2.3), Cm(30.8), Cm(10.6), [
        {"text": "1. 优先检查小麦 KL 的土层设定与实测土壤水分消耗是否一致。", "size": 16},
        {"text": "2. 对 DiffusSlope、DiffusConst、SWCON 做合理性边界检查，避免不物理的水分再分配。", "size": 16},
        {"text": "3. 用独立年份或处理验证 grain_yield、biomass、LAI、WUE 和土壤水分动态。", "size": 16},
        {"text": "4. 若论文篇幅允许，可把管理措施参数作为单独附加敏感性分析。", "size": 16},
        {"text": "5. 结果表述中区分品种参数不确定性和系统过程参数不确定性。", "size": 16},
    ], fill=RGBColor(255, 255, 255))

    slide = blank_slide(prs)
    add_title(slide, "Take-home messages")
    add_textbox(slide, Cm(1.4), Cm(2.5), Cm(30.4), Cm(11.4), [
        {"text": "1. 品种 Sobol 已经识别并固化主要物候和 RUE 参数。", "size": 17},
        {"text": "2. 系统 Sobol 显示根系取水和 SoilWater 参数是校准后模型的主要不确定来源。", "size": 17},
        {"text": "3. 小麦 KL 是当前最值得优先校准和独立验证的系统参数。", "size": 17},
        {"text": "4. N128 相比 N64 的 Top 参数稳定性较好，可作为正式结果。", "size": 17},
        {"text": "5. 后续工作应围绕根系取水、水分再分配和径流排水过程展开。", "size": 17},
    ], fill=RGBColor(255, 255, 255))

    ppt_path = OUT_DIR / "APSIM_Sobol_integrated_report_CN_fixed.pptx"
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(ppt_path)
    return ppt_path


def extract_text(ppt_path):
    prs = Presentation(str(ppt_path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)


def write_qa(ppt_path):
    prs = Presentation(str(ppt_path))
    with zipfile.ZipFile(ppt_path) as z:
        slides = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        media = [n for n in z.namelist() if n.startswith("ppt/media/")]
    text = extract_text(ppt_path)
    has_bad_question_marks = bool(re.search(r"\?{2,}", text))
    qa_path = OUT_DIR / "qa_report.md"
    qa_path.write_text(
        "\n".join(
            [
                "# QA report",
                "",
                f"- PPT: `{ppt_path.name}`",
                f"- slides opened by python-pptx: {len(prs.slides)}",
                f"- slide xml files in package: {len(slides)}",
                f"- media files in package: {len(media)}",
                f"- contains repeated question marks: {has_bad_question_marks}",
                f"- Chinese font requested: {FONT_NAME}",
                f"- matplotlib font file: `{FONT_PATH}`",
                "",
                "## Text preview",
                "",
                text[:1200],
            ]
        ),
        encoding="utf-8",
    )
    if has_bad_question_marks:
        raise RuntimeError("PPT 文本中仍然存在连续问号，请检查编码。")
    return qa_path, len(prs.slides), len(media)


def main():
    ensure_inputs()
    sys_df, cul_df, stab_df, best, metrics = prepare_data()
    figs = build_figures(sys_df, cul_df, stab_df)
    ppt_path = build_ppt(sys_df, cul_df, stab_df, best, metrics, figs)
    qa_path, slide_count, media_count = write_qa(ppt_path)
    print(f"PPT={ppt_path}")
    print(f"QA={qa_path}")
    print(f"slides={slide_count}")
    print(f"media={media_count}")


if __name__ == "__main__":
    main()
