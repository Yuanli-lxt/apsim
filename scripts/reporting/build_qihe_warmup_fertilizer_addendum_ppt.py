"""Build a Chinese PPTX addendum for Qihe fertilizer and N warm-up results."""

from __future__ import annotations

import argparse
import json
import zipfile
from pathlib import Path

import matplotlib as mpl
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
REPORT = ROOT / "report" / "2026_7_15"
ASSETS = REPORT / "assets" / "figures"
OUT = REPORT / "qihe_fertilizer_warmup_sensitivity_addendum_cn.pptx"
QA = REPORT / "warmup_addendum_qa_report.md"
MANIFEST = REPORT / "warmup_addendum_asset_manifest.md"
RUN = ROOT / "outputs" / "spatial" / "county_pilot_2020" / "warmup_sensitivity" / "qihe_warmup_10km_full_2010_2023_20260716_v1"
CONFIG = ROOT / "configs" / "spatial" / "qihe_warmup_sensitivity_scenarios.json"
OFFICIAL_ANNUAL = RUN / "annual_n_scenarios_vs_official.csv"
OFFICIAL_METRICS = RUN / "n_scenario_official_metrics.csv"

NAVY = "17324D"; TEAL = "2D7F82"; BLUE = "4B83B6"; GOLD = "C49A42"
RED = "B95756"; INK = "1F252B"; GREY = "66717B"; LIGHT = "F4F6F7"
PALE_TEAL = "DCECEA"; PALE_BLUE = "DDEAF2"; WHITE = "FFFFFF"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def mpl_color(value: str) -> str:
    return f"#{value}"


def add_text(slide, text, x, y, w, h, size=15, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame; frame.clear(); frame.word_wrap = True; frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]; paragraph.alignment = align
    run = paragraph.add_run(); run.text = text
    run.font.name = "Microsoft YaHei"; run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_bullets(slide, items, x, y, w, h, size=14, spacing=8, color=INK):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame; frame.clear(); frame.word_wrap = True
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = f"• {item}"; p.font.name = "Microsoft YaHei"; p.font.size = Pt(size)
        p.font.color.rgb = rgb(color); p.space_after = Pt(spacing); p.line_spacing = 1.12
    return shape


def add_rect(slide, x, y, w, h, fill=LIGHT, line="E0E5E8", rounded=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill); shape.line.color.rgb = rgb(line)
    return shape


def add_picture_contain(slide, path: Path, x, y, w, h):
    with Image.open(path) as image:
        ratio = image.width / image.height
    box = w / h
    if ratio >= box:
        pw, ph = w, w / ratio; px, py = x, y + (h - ph) / 2
    else:
        ph, pw = h, h * ratio; py, px = y, x + (w - pw) / 2
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def title(slide, text, number):
    add_text(slide, text, 0.55, 0.28, 11.9, 0.55, 25, NAVY, True)
    add_text(slide, f"{number:02d}", 12.25, 0.30, 0.48, 0.34, 10, GREY, True, PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.56), Inches(0.91), Inches(12.15), Inches(0.025))
    line.fill.solid(); line.fill.fore_color.rgb = rgb("D8DEE2"); line.line.fill.background()


def footer(slide, source="齐河县APSIM连续模拟补充分析 · 2026-07-16"):
    add_text(slide, source, 0.58, 7.10, 11.7, 0.20, 7.5, GREY)
    add_text(slide, "APSIM Classic 7.10 r4221", 11.05, 7.10, 1.65, 0.20, 7.5, GREY, align=PP_ALIGN.RIGHT)


def takeaway(slide, text):
    add_rect(slide, 0.65, 6.43, 12.0, 0.46, PALE_TEAL, line=PALE_TEAL, rounded=False)
    add_text(slide, text, 0.88, 6.52, 11.55, 0.24, 11.5, TEAL, True, PP_ALIGN.CENTER)


def metric(slide, value, label, x, y, w=2.65, color=TEAL):
    add_rect(slide, x, y, w, 1.15, LIGHT)
    add_text(slide, value, x + 0.10, y + 0.14, w - 0.20, 0.42, 23, color, True, PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.10, y + 0.70, w - 0.20, 0.24, 9.5, GREY, align=PP_ALIGN.CENTER)


def add_table(slide, rows, x, y, w, h, widths=None):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    if widths:
        for col, width in zip(table.columns, widths): col.width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c); cell.text = str(value); cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
            cell.fill.solid(); cell.fill.fore_color.rgb = rgb(NAVY if r == 0 else (WHITE if r % 2 else LIGHT))
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"; run.font.size = Pt(10 if r else 10.5)
                    run.font.bold = r == 0; run.font.color.rgb = rgb(WHITE if r == 0 else INK)
    return table


def load_data():
    annual = pd.read_csv(RUN / "annual_county_yield_and_mineral_n.csv")
    conv = pd.read_csv(RUN / "warmup_convergence_by_year.csv")
    official_annual = pd.read_csv(OFFICIAL_ANNUAL)
    official_metrics = pd.read_csv(OFFICIAL_METRICS)
    config = json.loads(CONFIG.read_text(encoding="utf-8"))
    return annual, conv, official_annual, official_metrics, config


def build_assets(annual: pd.DataFrame, conv: pd.DataFrame, official_annual: pd.DataFrame,
                 official_metrics: pd.DataFrame, config: dict) -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    mpl.rcParams.update({
        "font.family": "sans-serif", "font.sans-serif": ["Microsoft YaHei", "SimHei", "Arial", "DejaVu Sans"],
        "font.size": 9, "axes.spines.top": False, "axes.spines.right": False, "axes.linewidth": 0.8,
    })
    formal = annual.query("year >= 2015 and initial_n_multiplier == 1").copy()

    convergence = ASSETS / "warmup_convergence_2011_2023.png"
    worst = conv.groupby("year").agg(
        wheat=("wheat_yield_relative_range_percent", "max"),
        maize=("maize_yield_relative_range_percent", "max"),
        mineral_n=("oct1_mineral_n_relative_range_percent", "max"),
    )
    fig, axes = plt.subplots(1, 2, figsize=(9.0, 3.45))
    axes[0].plot(worst.index, worst.wheat, "o-", color=mpl_color(BLUE), lw=1.8, label="小麦")
    axes[0].plot(worst.index, worst.maize, "o-", color=mpl_color(TEAL), lw=1.8, label="玉米")
    axes[0].axhline(2, color=mpl_color(RED), ls="--", lw=1.2, label="阈值 2%")
    axes[0].axvspan(2015, 2023.4, color=mpl_color(PALE_TEAL), alpha=0.65)
    axes[0].set(title="产量对初始氮的最大相对极差", ylabel="相对极差 (%)", xticks=[2011, 2013, 2015, 2017, 2019, 2021, 2023])
    axes[0].legend(frameon=False, ncol=3, fontsize=8)
    axes[1].plot(worst.index, worst.mineral_n, "o-", color=mpl_color(GOLD), lw=1.9)
    axes[1].axhline(10, color=mpl_color(RED), ls="--", lw=1.2, label="阈值 10%")
    axes[1].axvspan(2015, 2023.4, color=mpl_color(PALE_TEAL), alpha=0.65)
    axes[1].set(title="10月1日矿质氮的最大相对极差", ylabel="相对极差 (%)", xticks=[2011, 2013, 2015, 2017, 2019, 2021, 2023])
    axes[1].legend(frameon=False, fontsize=8)
    for ax in axes: ax.set_xlim(2010.6, 2023.4); ax.grid(axis="y", alpha=0.18)
    fig.tight_layout(); fig.savefig(convergence, dpi=280, bbox_inches="tight", facecolor="white"); plt.close(fig)

    means = ASSETS / "warmup_scenario_mean_yield_mineral_n.png"
    order = ["fixed_180_reference", "statistical_constraint_central", "statistical_constraint_low", "statistical_constraint_high"]
    labels = ["固定180", "统计中央", "中央×0.75", "中央×1.25"]
    summary = formal.groupby("management_scenario")[["wheat_yield_kg_ha", "maize_yield_kg_ha", "oct1_mineral_n_kg_ha"]].mean().loc[order]
    fig, axes = plt.subplots(1, 2, figsize=(9.1, 3.5), gridspec_kw={"width_ratios": [1.45, 1]})
    x = np.arange(4)
    axes[0].bar(x - 0.18, summary.wheat_yield_kg_ha, 0.36, color=mpl_color(BLUE), label="小麦")
    axes[0].bar(x + 0.18, summary.maize_yield_kg_ha, 0.36, color=mpl_color(TEAL), label="玉米")
    axes[0].set(xticks=x, xticklabels=labels, ylabel="产量 (kg/ha)", title="2015—2023平均原始APSIM产量")
    axes[0].tick_params(axis="x", labelrotation=12); axes[0].legend(frameon=False, ncol=2)
    axes[1].bar(x, summary.oct1_mineral_n_kg_ha, 0.62, color=[mpl_color(GREY), mpl_color(GOLD), mpl_color(BLUE), mpl_color(RED)])
    axes[1].set(xticks=x, xticklabels=labels, ylabel="NO3+NH4 (kg/ha)", title="平均播前矿质氮")
    axes[1].tick_params(axis="x", labelrotation=12)
    for ax in axes: ax.grid(axis="y", alpha=0.18)
    fig.tight_layout(); fig.savefig(means, dpi=280, bbox_inches="tight", facecolor="white"); plt.close(fig)

    official_error = ASSETS / "warmup_12_scenarios_official_error.png"
    scenario_labels = {
        "fixed_180_reference": "固定180",
        "statistical_constraint_central": "统计中央",
        "statistical_constraint_low": "中央×0.75",
        "statistical_constraint_high": "中央×1.25",
    }
    fig, axes = plt.subplots(1, 2, figsize=(9.25, 3.65))
    for ax, crop, crop_label in zip(axes, ["wheat", "maize"], ["小麦", "玉米"]):
        d = official_metrics.query("crop == @crop").copy()
        matrix = d.pivot(index="management_scenario", columns="initial_n_multiplier", values="mape_percent").loc[order, [0.5, 1.0, 1.5]]
        im = ax.imshow(matrix, cmap="YlOrRd", vmin=0, vmax=70, aspect="auto")
        ax.set(xticks=np.arange(3), xticklabels=["0.5×", "1.0×", "1.5×"],
               yticks=np.arange(4), yticklabels=[scenario_labels[x] for x in order],
               xlabel="初始NO3/NH4倍数", title=f"{crop_label} MAPE (%)")
        for i, scenario in enumerate(order):
            for j, multiplier in enumerate([0.5, 1.0, 1.5]):
                row = d[(d.management_scenario == scenario) & (d.initial_n_multiplier == multiplier)].iloc[0]
                color = "white" if row.mape_percent > 40 else mpl_color(INK)
                ax.text(j, i, f"{row.mape_percent:.1f}%\n偏差 {row.mean_bias_kg_ha/1000:+.2f}",
                        ha="center", va="center", fontsize=7.4, color=color)
    fig.text(0.5, 0.02, "单元格：MAPE；下行：平均偏差（t/ha）", ha="center", fontsize=8, color=mpl_color(GREY))
    fig.subplots_adjust(left=0.14, right=0.99, bottom=0.19, top=0.86, wspace=0.38)
    fig.savefig(official_error, dpi=280, bbox_inches="tight", facecolor="white"); plt.close(fig)

    official_envelope = ASSETS / "warmup_12_scenarios_official_annual_envelope.png"
    fig, axes = plt.subplots(1, 2, figsize=(9.25, 3.55), sharex=True)
    for ax, crop, crop_label in zip(axes, ["wheat", "maize"], ["小麦", "玉米"]):
        d = official_annual.query("crop == @crop")
        envelope = d.groupby("year").agg(
            official=("official_yield_kg_ha", "first"),
            minimum=("raw_apsim_yield_kg_ha", "min"),
            maximum=("raw_apsim_yield_kg_ha", "max"),
        )
        central = d.query("management_scenario == 'statistical_constraint_central' and initial_n_multiplier == 1.0").set_index("year")
        ax.fill_between(envelope.index, envelope.minimum, envelope.maximum, color=mpl_color(PALE_TEAL), alpha=0.85, label="12种组合范围")
        ax.plot(central.index, central.raw_apsim_yield_kg_ha, "o-", color=mpl_color(TEAL), lw=1.7, ms=3.5, label="统计中央/初始N=1.0")
        ax.plot(envelope.index, envelope.official, "o-", color=mpl_color(INK), lw=2.1, ms=4.0, label="正式统计")
        ax.set(title=crop_label, ylabel="单产 (kg/ha)", xticks=[2015, 2017, 2019, 2021, 2023])
        ax.grid(axis="y", alpha=0.18)
    axes[0].legend(frameon=False, fontsize=7.6, ncol=1, loc="lower left")
    fig.tight_layout(); fig.savefig(official_envelope, dpi=280, bbox_inches="tight", facecolor="white"); plt.close(fig)

    legacy = ASSETS / "warmup_central_n_rate_and_legacy.png"
    central = formal.query("management_scenario == 'statistical_constraint_central'").sort_values("year")
    rates = {int(k): float(v) for k, v in config["management_scenarios"]["statistical_constraint_central"]["annual_total_n_kg_ha"].items()}
    central["n_rate"] = central.year.map(rates)
    fig, ax = plt.subplots(figsize=(8.3, 3.55)); ax2 = ax.twinx()
    ax.plot(central.year, central.n_rate, "o-", color=mpl_color(TEAL), lw=2.1, label="统计中央施氮")
    ax2.plot(central.year, central.oct1_mineral_n_kg_ha, "s-", color=mpl_color(GOLD), lw=2.1, label="10月1日矿质氮")
    ax.axhline(180, color=mpl_color(GREY), ls="--", lw=1.1, label="固定180")
    ax.set(xlabel="年份", ylabel="施氮强度 (kg N/ha)", xticks=central.year); ax2.set_ylabel("矿质氮 (kg/ha)")
    lines = ax.get_lines() + ax2.get_lines(); ax.legend(lines, [line.get_label() for line in lines], frameon=False, ncol=3, loc="upper right")
    ax.grid(axis="y", alpha=0.18); fig.tight_layout(); fig.savefig(legacy, dpi=280, bbox_inches="tight", facecolor="white"); plt.close(fig)

    yields = ASSETS / "warmup_annual_yield_scenarios.png"
    fig, axes = plt.subplots(1, 2, figsize=(9.2, 3.5), sharex=True)
    colors = [GREY, TEAL, BLUE, RED]
    for ax, crop, crop_label in zip(axes, ["wheat", "maize"], ["小麦", "玉米"]):
        for scenario, label, color in zip(order, labels, colors):
            d = formal[formal.management_scenario == scenario].sort_values("year")
            ax.plot(d.year, d[f"{crop}_yield_kg_ha"], "o-", lw=1.6, ms=3.5, color=mpl_color(color), label=label)
        ax.set(title=crop_label, ylabel="原始APSIM产量 (kg/ha)", xticks=[2015, 2017, 2019, 2021, 2023]); ax.grid(axis="y", alpha=0.18)
    axes[0].legend(frameon=False, fontsize=7.5, ncol=2)
    fig.tight_layout(); fig.savefig(yields, dpi=280, bbox_inches="tight", facecolor="white"); plt.close(fig)

    response = ASSETS / "warmup_low_high_response_heatmap.png"
    pivot = formal.pivot(index="year", columns="management_scenario", values=["wheat_yield_kg_ha", "maize_yield_kg_ha"])
    matrix=[]; row_labels=[]
    for crop, crop_label in [("wheat", "小麦"), ("maize", "玉米")]:
        variable=f"{crop}_yield_kg_ha"; central_y=pivot[variable]["statistical_constraint_central"]
        for scenario, label in [("statistical_constraint_low", "低—中央"), ("statistical_constraint_high", "高—中央")]:
            matrix.append(((pivot[variable][scenario] / central_y - 1) * 100).to_numpy()); row_labels.append(f"{crop_label} {label}")
    matrix=np.asarray(matrix)
    fig, ax = plt.subplots(figsize=(9.0, 2.8)); im=ax.imshow(matrix, cmap="RdBu_r", vmin=-30, vmax=30, aspect="auto")
    ax.set(xticks=np.arange(9), xticklabels=range(2015,2024), yticks=np.arange(4), yticklabels=row_labels)
    for i in range(4):
        for j in range(9):
            value=matrix[i,j]; ax.text(j,i,f"{value:.0f}%",ha="center",va="center",fontsize=7.5,color="white" if abs(value)>16 else mpl_color(INK))
    fig.colorbar(im, ax=ax, pad=0.015, label="相对统计中央产量变化 (%)")
    ax.set_title("低/高施氮相对统计中央值的逐年产量响应")
    fig.tight_layout(); fig.savefig(response, dpi=280, bbox_inches="tight", facecolor="white"); plt.close(fig)
    return {"convergence": convergence, "means": means, "official_error": official_error,
            "official_envelope": official_envelope, "legacy": legacy, "yields": yields, "response": response}


def build_deck(assets: dict[str, Path]) -> Presentation:
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide=prs.slides.add_slide(blank)
    add_rect(slide,0,0,13.333,7.5,WHITE,WHITE,False); add_rect(slide,0,0,0.24,7.5,TEAL,TEAL,False)
    add_text(slide,"施肥统计约束与连续氮预热",0.85,1.25,11.6,0.72,29,NAVY,True)
    add_text(slide,"齐河县冬小麦—夏玉米APSIM补充分析",0.88,2.08,10.8,0.48,19,TEAL,True)
    add_text(slide,"仅补充昨日汇报未覆盖的内容：化肥统计折纯与分配假设 · 2010—2023连续模拟 · 初始氮收敛 · 12种氮情景与正式统计对照",0.90,2.83,10.9,0.65,13,GREY)
    metric(slide,"1056/1056","连续APSIM案例成功",0.9,4.25,2.75,TEAL); metric(slide,"2010—2023","真实逐日AgERA5",3.9,4.25,2.75,BLUE)
    metric(slide,"4 × 3","施氮管理 × 初始氮",6.9,4.25,2.75,GOLD); metric(slide,"10 km","县域批量诊断",9.9,4.25,2.55,RED)
    add_text(slide,"2026-07-16 · APSIM Classic 7.10 r4221",0.9,6.72,5.5,0.28,9,GREY)
    add_notes(slide,"本PPT是昨日县域APSIM工作总结的补充专题，不重复空间掩膜、固定产量校准和分辨率分析。重点回答统计施氮中央值是什么、预热能否降低初始氮影响、不同氮投入如何改变产量和土壤矿质氮，以及12种组合与2015至2023年正式统计单产相比表现如何。")

    slide=prs.slides.add_slide(blank); title(slide,"统计中央值是县级化肥总量约束下的参考施氮强度",2)
    add_rect(slide,0.72,1.22,7.35,2.0,PALE_BLUE,line=PALE_BLUE)
    add_text(slide,"中央施氮 = [折纯单质氮 + 复合肥折纯量 × 1/3] ÷ 农作物总播种面积",1.03,1.63,6.72,0.72,18,NAVY,True,PP_ALIGN.CENTER)
    add_text(slide,"1/3代表平衡型复合肥N–P₂O₅–K₂O≈1:1:1假设；按总播种面积分配代表每作物季公顷。",1.10,2.52,6.60,0.35,10.5,GREY,align=PP_ALIGN.CENTER)
    add_rect(slide,8.45,1.22,4.18,4.55,LIGHT)
    add_text(slide,"2020年计算示例",8.78,1.55,3.50,0.38,16,NAVY,True)
    add_bullets(slide,["折纯单质氮：20,125 t","复合肥折纯量：20,844 t","估算氮总量：27,073 t N","总播种面积：170,800 ha"],8.82,2.13,3.42,1.85,12,7)
    add_text(slide,"158.51 kg N/ha",8.78,4.55,3.52,0.55,23,TEAL,True,PP_ALIGN.CENTER)
    add_bullets(slide,["官方量：县级、全部农业用途","假设量：复合肥N比例与作物分配","不是小麦/玉米实测施氮量"],0.92,3.70,6.72,1.48,13,8)
    takeaway(slide,"“中央”表示不确定分配中的参考情景，不表示官方推荐量或真实田间均值。")
    footer(slide,"来源：德州统计年鉴2013—2024；配置：qihe_warmup_sensitivity_scenarios.json")
    add_notes(slide,"官方年鉴直接给出的是县级化肥折纯量和农作物播种面积。中央值还包含两个假设：复合肥折纯养分的三分之一计为氮，以及全部作物按播种面积均匀分配。因此它是统计约束情景，而不是观测到的作物季施氮。")

    slide=prs.slides.add_slide(blank); title(slide,"连续试验将管理不确定性与初始氮不确定性正交分开",3)
    labels=[("真实逐日气象","AgERA5 2010—2023\n25个节点",BLUE),("空间土壤","14个HWSD单元\n88个独立案例",TEAL),("施氮管理","固定180 / 中央\n中央×0.75 / ×1.25",GOLD),("初始矿质氮","NO₃与NH₄同时\n×0.5 / ×1.0 / ×1.5",RED)]
    for i,(head,body,color) in enumerate(labels):
        x=0.68+i*3.14; add_rect(slide,x,1.35,2.78,1.62,LIGHT,line=color); add_text(slide,head,x+0.16,1.61,2.46,0.32,14,color,True,PP_ALIGN.CENTER); add_text(slide,body,x+0.17,2.10,2.44,0.58,11,INK,align=PP_ALIGN.CENTER)
        if i<3: add_text(slide,"×",x+2.87,1.88,0.27,0.38,19,GREY,True,PP_ALIGN.CENTER)
    add_rect(slide,0.80,3.50,11.76,1.55,PALE_TEAL,line=PALE_TEAL)
    add_text(slide,"2010-10",1.05,3.78,1.2,0.3,12,NAVY,True); add_text(slide,"预热期：2010—2014",2.25,3.78,3.20,0.3,14,GOLD,True,PP_ALIGN.CENTER); add_text(slide,"正式分析：2015—2023",7.15,3.78,3.35,0.3,14,TEAL,True,PP_ALIGN.CENTER)
    line=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(1.15),Inches(4.42),Inches(10.70),Inches(0.06)); line.fill.solid(); line.fill.fore_color.rgb=rgb(TEAL); line.line.fill.background()
    add_text(slide,"土壤水、矿质氮、有机质和残茬逐年连续传递，不做年度重置",2.05,5.43,9.2,0.36,14,NAVY,True,PP_ALIGN.CENTER)
    takeaway(slide,"改变的只有年度施氮方案与初始NO₃/NH₄倍数；品种、生理参数和其他管理保持不变。")
    footer(slide); add_notes(slide,"这是一个正交敏感性设计。不同土壤和天气节点保留真实差异，三个初始氮水平用于检验预热；四个管理成员用于检验年度投入。2010至2011施氮缺少在线年鉴，明确沿用2012强度作为预热填补。")

    slide=prs.slides.add_slide(blank); title(slide,"到2015年，初始氮对产量和播前矿质氮的影响降至预设阈值内",4)
    add_picture_contain(slide,assets["convergence"],0.60,1.18,9.15,4.95)
    add_rect(slide,9.95,1.42,2.72,4.38,LIGHT)
    add_text(slide,"正式期最大值",10.20,1.72,2.18,0.30,14,NAVY,True,PP_ALIGN.CENTER)
    add_text(slide,"0.53%",10.22,2.35,2.12,0.43,23,BLUE,True,PP_ALIGN.CENTER); add_text(slide,"小麦产量",10.22,2.80,2.12,0.24,9,GREY,align=PP_ALIGN.CENTER)
    add_text(slide,"0.23%",10.22,3.37,2.12,0.43,23,TEAL,True,PP_ALIGN.CENTER); add_text(slide,"玉米产量",10.22,3.82,2.12,0.24,9,GREY,align=PP_ALIGN.CENTER)
    add_text(slide,"2.90%",10.22,4.39,2.12,0.43,23,GOLD,True,PP_ALIGN.CENTER); add_text(slide,"10月1日矿质氮",10.22,4.84,2.12,0.24,9,GREY,align=PP_ALIGN.CENTER)
    takeaway(slide,"预热使正式期结果不再由任意初始氮控制，但不等于初始氮本身经过实测验证。")
    footer(slide,"来源：warmup_convergence_by_year.csv；阴影为2015—2023正式分析期")
    add_notes(slide,"2011至2014并非全部通过，尤其2011差异最大，2014统计中央情景的小麦仍有一次3.60%的未通过。结论应限定为从2015年起持续低于操作阈值，不能说初始状态被物理消除。")

    slide=prs.slides.add_slide(blank); title(slide,"施氮方案差异远大于预热后的初始氮差异",5)
    add_picture_contain(slide,assets["means"],0.52,1.15,9.35,5.0)
    add_rect(slide,9.98,1.35,2.72,4.55,LIGHT)
    add_text(slide,"相对统计中央",10.22,1.66,2.20,0.32,14,NAVY,True,PP_ALIGN.CENTER)
    add_text(slide,"低情景",10.22,2.23,2.20,0.26,11,BLUE,True,PP_ALIGN.CENTER)
    add_text(slide,"小麦 −12.7%\n玉米 −10.5%",10.22,2.58,2.20,0.72,16,INK,True,PP_ALIGN.CENTER)
    add_text(slide,"高情景",10.22,3.58,2.20,0.26,11,RED,True,PP_ALIGN.CENTER)
    add_text(slide,"小麦 +5.6%\n玉米 +6.9%",10.22,3.93,2.20,0.72,16,INK,True,PP_ALIGN.CENTER)
    add_text(slide,"矿质氮 +106%",10.22,5.02,2.20,0.34,14,GOLD,True,PP_ALIGN.CENTER)
    takeaway(slide,"高投入的氮残留增幅远大于产量增幅，提示潜在环境代价。")
    footer(slide,"2015—2023县域面积加权均值；1.0倍初始氮；产量未经统计乘法校正")
    add_notes(slide,"中央值高情景平均增产约6%，但播前矿质氮翻倍，说明额外投入并未等比例转化为籽粒。低情景产量下降更明显，说明后期可能发生氮限制。这里的产量是原始APSIM输出，不宜直接作为县域推荐施氮量。")

    slide=prs.slides.add_slide(blank); title(slide,"与正式统计比较后，12种组合仍有作物方向一致的系统偏差",6)
    add_picture_contain(slide,assets["official_error"],0.52,1.10,9.35,5.05)
    add_rect(slide,9.98,1.36,2.72,4.56,LIGHT)
    add_text(slide,"最低MAPE组合",10.20,1.68,2.25,0.30,14,NAVY,True,PP_ALIGN.CENTER)
    add_text(slide,"12.3%",10.20,2.18,2.25,0.42,23,BLUE,True,PP_ALIGN.CENTER)
    add_text(slide,"小麦：固定180 / 初始N 1.5×\n平均偏差 −0.94 t/ha",10.16,2.64,2.34,0.66,10.2,GREY,align=PP_ALIGN.CENTER)
    add_text(slide,"38.0%",10.20,3.62,2.25,0.42,23,TEAL,True,PP_ALIGN.CENTER)
    add_text(slide,"玉米：中央×0.75 / 初始N 0.5×\n平均偏差 +2.98 t/ha",10.16,4.08,2.34,0.66,10.2,GREY,align=PP_ALIGN.CENTER)
    add_text(slide,"同一管理内，初始N倍数几乎不改变误差",10.18,5.15,2.30,0.46,11,RED,True,PP_ALIGN.CENTER)
    takeaway(slide,"改变氮投入可以移动产量量级，但不能同时修正小麦低估与玉米高估。")
    footer(slide,"正式统计：齐河县2015—2023县级单产；原始APSIM未施加产量校正系数")
    add_notes(slide,"这里补上正式统计对照。每个单元格报告2015至2023年的MAPE和平均偏差，12种组合均未针对正式统计重新拟合。所谓最低MAPE仅是描述性排序，不是独立验证或推荐方案。小麦所有组合整体低估，玉米所有组合整体高估，说明当前误差不能只归因于初始氮或施氮强度。")

    slide=prs.slides.add_slide(blank); title(slide,"12种施氮—初始氮组合不能稳定覆盖官方年际变化",7)
    add_picture_contain(slide,assets["official_envelope"],0.52,1.12,9.45,5.00)
    add_rect(slide,10.08,1.38,2.62,4.52,LIGHT)
    add_text(slide,"官方值落入情景范围",10.24,1.70,2.28,0.32,13,NAVY,True,PP_ALIGN.CENTER)
    add_text(slide,"3 / 9年",10.24,2.28,2.28,0.45,23,BLUE,True,PP_ALIGN.CENTER)
    add_text(slide,"小麦",10.24,2.76,2.28,0.25,10,GREY,align=PP_ALIGN.CENTER)
    add_text(slide,"0 / 9年",10.24,3.55,2.28,0.45,23,TEAL,True,PP_ALIGN.CENTER)
    add_text(slide,"玉米：所有年份官方值均低于模拟下界",10.20,4.03,2.36,0.52,10,GREY,align=PP_ALIGN.CENTER)
    add_text(slide,"年际相关性普遍较弱，部分组合为负",10.20,5.08,2.36,0.50,11,RED,True,PP_ALIGN.CENTER)
    takeaway(slide,"氮情景主要改变平均量级，并未恢复与正式统计一致的年际响应。")
    footer(slide,"阴影：每年12种原始APSIM组合的最小—最大范围；黑线：正式统计单产")
    add_notes(slide,"阴影不是统计置信区间，而是12种离散组合的范围。小麦只有2018、2020和2022年官方值落入该范围；玉米九年全部低于模拟范围。相关性也没有显示稳定改善，因此不能以扩大氮情景范围代替对水分、管理、物候、面积口径和模型结构的诊断。")

    slide=prs.slides.add_slide(blank); title(slide,"统计中央情景经历“前期氮库积累—后期快速消耗”",8)
    add_picture_contain(slide,assets["legacy"],0.62,1.18,8.65,4.9)
    add_rect(slide,9.48,1.40,3.17,4.43,LIGHT)
    add_text(slide,"关键转折",9.80,1.73,2.52,0.32,15,NAVY,True)
    add_bullets(slide,["2015：260.8 kg N/ha","2015播前矿质氮：376.7 kg/ha","2021后中央施氮低于140 kg N/ha","2022播前矿质氮降至29.4 kg/ha","2023小幅回升至42.7 kg/ha"],9.78,2.25,2.55,2.45,11.5,7)
    add_text(slide,"年际产量不能只用当年施肥解释",9.74,5.10,2.62,0.46,13,RED,True,PP_ALIGN.CENTER)
    takeaway(slide,"连续轮作保留了显著氮遗留效应；前期盈余可暂时缓冲后期低投入。")
    footer(slide,"统计中央值：2012—2023年鉴约束；2010—2011使用2012值进行预热")
    add_notes(slide,"统计中央施氮逐年下降，但矿质氮并未同步下降，因为前期投入形成了土壤氮库。到2021至2022残余氮快速降低，低施氮影响开始放大。矿质氮绝对量尚无田间观测验证，主要用于比较情景和趋势。")

    slide=prs.slides.add_slide(blank); title(slide,"玉米在后期表现出更清晰的氮限制，小麦响应受连续过程耦合",9)
    add_picture_contain(slide,assets["yields"],0.58,1.15,9.20,4.95)
    add_rect(slide,9.98,1.38,2.72,4.55,LIGHT)
    add_text(slide,"玉米",10.22,1.70,2.20,0.30,14,TEAL,True,PP_ALIGN.CENTER)
    add_bullets(slide,["2020—2022低情景下降19%—22%","同期高情景增加12%—15%","2023高投入几乎不再增产"],10.22,2.12,2.18,1.52,10.8,7)
    add_text(slide,"小麦",10.22,3.92,2.20,0.30,14,BLUE,True,PP_ALIGN.CENTER)
    add_bullets(slide,["仅4/9年呈高≥中央≥低","2022高情景比中央高41%","需检查水分、物候与前茬效应"],10.22,4.32,2.18,1.34,10.8,7)
    takeaway(slide,"玉米施氮响应较稳定；小麦不能用简单的单季剂量—产量关系解释。")
    footer(slide,"来源：annual_county_yield_and_mineral_n.csv；所有情景使用相同气象、土壤、品种与基础管理")
    add_notes(slide,"玉米大多数年份满足随氮增加而增产，后期低氮损失尤其明显。小麦在多个年份出现非单调响应，可能与前茬耗水耗氮、残茬矿化、物候或强制收获有关，当前结果不足以确定机制。")

    slide=prs.slides.add_slide(blank); title(slide,"逐年响应显示：高氮并非每年有效，低氮风险在后期集中暴露",10)
    add_picture_contain(slide,assets["response"],0.62,1.25,9.55,4.5)
    add_rect(slide,10.30,1.45,2.34,4.20,LIGHT)
    add_text(slide,"读图规则",10.55,1.78,1.82,0.30,14,NAVY,True,PP_ALIGN.CENTER)
    add_text(slide,"红：高于中央\n蓝：低于中央",10.54,2.27,1.84,0.68,13,INK,True,PP_ALIGN.CENTER)
    add_text(slide,"最强低氮损失",10.53,3.38,1.86,0.28,11,GREY,align=PP_ALIGN.CENTER)
    add_text(slide,"2022小麦 −27%\n2022玉米 −22%",10.52,3.77,1.88,0.72,15,BLUE,True,PP_ALIGN.CENTER)
    add_text(slide,"高氮异常",10.53,4.76,1.86,0.28,11,GREY,align=PP_ALIGN.CENTER)
    add_text(slide,"2023玉米 −0.2%",10.52,5.10,1.88,0.30,13,RED,True,PP_ALIGN.CENTER)
    takeaway(slide,"同一施氮调整在不同年份产生不同回报，固定比例敏感性不是施肥推荐曲线。")
    footer(slide,"色块为相对统计中央情景的产量变化；比较限定在同一年、同一初始氮水平")
    add_notes(slide,"热图用于强调年际异质性。2022对低氮最敏感，而2023高氮对玉米几乎没有收益，说明天气、水分或其他限制可改变氮肥边际效应。不能把0.75和1.25两个点拟合成通用推荐曲线。")

    slide=prs.slides.add_slide(blank); title(slide,"结论：氮预热已收敛，但氮设置不能解释产量真值偏差",11)
    cards=[("可以确认",["2015后初始氮影响低于预设阈值","12种组合已与2015—2023统计对照","管理差异大于初始状态差异"],TEAL),("不能确认",["最低误差组合不是独立验证","氮强度不能同时修正两作物偏差","原始产量不能直接形成推荐量"],RED),("下一步优先",["保留氮情景作为不确定性范围","诊断水分、物候、面积与管理过程","用土壤氮观测和独立年份再验证"],GOLD)]
    for i,(head,items,color) in enumerate(cards):
        x=0.70+i*4.18; add_rect(slide,x,1.35,3.78,4.55,LIGHT,line=color); add_text(slide,head,x+0.22,1.72,3.34,0.35,16,color,True,PP_ALIGN.CENTER); add_bullets(slide,items,x+0.30,2.35,3.18,2.25,12.2,12); add_text(slide,["证据充分","证据边界","验证路线"][i],x+0.42,5.22,2.94,0.32,11,GREY,True,PP_ALIGN.CENTER)
    takeaway(slide,"正式统计对照表明：不能通过选择某个氮情景替代独立校准、过程诊断和跨年验证。")
    footer(slide); add_notes(slide,"阶段性结论是初始值敏感性已经通过预热得到控制，但12种氮组合与正式统计比较后仍存在明显系统偏差。下一阶段不应从这12种组合中事后挑选所谓最佳方案，更不应逐年调氮拟合产量；应保留氮情景用于不确定性分析，并增加逐年施肥、土壤无机氮、水分胁迫、物候、面积口径和氮损失证据。")
    return prs


def validate_and_write(prs: Presentation, assets: dict[str, Path]):
    REPORT.mkdir(parents=True, exist_ok=True); prs.save(OUT)
    reopened=Presentation(OUT); names=zipfile.ZipFile(OUT).namelist(); media=[n for n in names if n.startswith("ppt/media/")]
    notes=sum(bool(slide.notes_slide.notes_text_frame.text.strip()) for slide in reopened.slides)
    overflow=[]
    for i,slide in enumerate(reopened.slides,1):
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > reopened.slide_width or shape.top + shape.height > reopened.slide_height:
                overflow.append(f"slide {i}: {shape.name}")
    QA.write_text(
        "# Warm-up addendum PPTX QA report\n\n"
        f"- PPTX creation: PASS\n- Output: `{OUT.name}`\n- Slide count: {len(reopened.slides)}\n"
        f"- Embedded media files: {len(media)}\n- Slides with speaker notes: {notes}\n"
        f"- Shape-bound violations: {len(overflow)}\n- Package verification: reopened with python-pptx and inspected as ZIP.\n"
        "- Existing 2026-07-15 PPTX and its assets were not modified.\n"
        "- Official comparison: 12 scenarios × 9 years × 2 crops; no scenario-specific yield recalibration.\n"
        "- Interpretation: descriptive external comparison, not independent validation or fertilizer recommendation.\n"
        "- Known limitation: no available headless PowerPoint/LibreOffice renderer; QA is structural, not pixel-rendered.\n"
        "- Manual follow-up: open once in PowerPoint to confirm Chinese font substitution.\n",
        encoding="utf-8",
    )
    rows=["# Warm-up addendum asset manifest","","| Asset | Source | Method | Slides |","|---|---|---|---|"]
    mapping={
        "convergence":("warmup_convergence_by_year.csv","matplotlib redraw without changing values","4"),
        "means":("annual_county_yield_and_mineral_n.csv","matplotlib aggregation; 2015—2023, initial N=1.0","5"),
        "official_error":("n_scenario_official_metrics.csv","MAPE and mean-bias heatmap; no refitting","6"),
        "official_envelope":("annual_n_scenarios_vs_official.csv","annual official line with 12-scenario min–max envelope","7"),
        "legacy":("annual_county_yield_and_mineral_n.csv + scenario config","matplotlib redraw","8"),
        "yields":("annual_county_yield_and_mineral_n.csv","matplotlib redraw","9"),
        "response":("annual_county_yield_and_mineral_n.csv","relative-difference heatmap","10"),
    }
    for key,path in assets.items():
        source,method,slides=mapping[key]; rows.append(f"| `{path.name}` | `{source}` | {method} | {slides} |")
    MANIFEST.write_text("\n".join(rows)+"\n",encoding="utf-8")
    return {"slides":len(reopened.slides),"media":len(media),"notes":notes,"overflow":overflow}


def main():
    global OUT
    parser=argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--force",action="store_true",help="Replace only the existing addendum PPTX and its dedicated QA files.")
    parser.add_argument("--output",type=Path,default=OUT,help="Output PPTX path; useful when the existing deck is open.")
    args=parser.parse_args()
    OUT=args.output.resolve()
    if OUT.exists() and not args.force:
        raise FileExistsError(f"Refusing to overwrite existing addendum: {OUT}")
    annual,conv,official_annual,official_metrics,config=load_data()
    assets=build_assets(annual,conv,official_annual,official_metrics,config)
    result=validate_and_write(build_deck(assets),assets)
    print(json.dumps({"output":str(OUT),**result},ensure_ascii=False,indent=2))


if __name__ == "__main__":
    main()
