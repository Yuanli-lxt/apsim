"""Build a Chinese PPTX summary of the recent Qihe county APSIM workflow."""

from __future__ import annotations

import json
import shutil
from pathlib import Path

import matplotlib as mpl
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
REPORT = ROOT / "report"
ASSETS = REPORT / "assets" / "figures"
OUT = REPORT / "qihe_county_APSIM_recent_work_summary_cn.pptx"
QA = REPORT / "qa_report.md"
MANIFEST = REPORT / "asset_manifest.md"

SPATIAL_OUT = ROOT / "outputs" / "spatial" / "county_pilot_2020"
VALIDATION = SPATIAL_OUT / "corrected_baseline" / "multiyear_resolution_validation"
STATS = ROOT / "data" / "processed" / "spatial" / "county_pilot_2020" / "calibration" / "qihe_2018_2020_official_crop_statistics.csv"

NAVY = "17324D"
TEAL = "2D7F82"
BLUE = "4B83B6"
PALE_BLUE = "DDEAF2"
PALE_TEAL = "DCECEA"
GOLD = "C49A42"
RED = "B95756"
INK = "1F252B"
GREY = "66717B"
LIGHT = "F4F6F7"
WHITE = "FFFFFF"


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def setup_assets() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    mpl.rcParams.update({
        "font.family": "sans-serif",
        "font.sans-serif": ["Microsoft YaHei", "SimHei", "Arial", "DejaVu Sans"],
        "font.size": 9,
        "axes.spines.top": False,
        "axes.spines.right": False,
        "axes.linewidth": 0.8,
        "pdf.fonttype": 42,
    })

    source_map = SPATIAL_OUT / "qihe_2020_rotation_grid_1km_5km_10km.png"
    source_qc = SPATIAL_OUT / "qihe_2020_crop_grid_quality_check.png"
    map_asset = ASSETS / "rotation_grid_1km_5km_10km.png"
    shutil.copy2(source_map, map_asset)

    qc_crop = ASSETS / "crop_masks_wheat_maize_rotation.png"
    with Image.open(source_qc) as image:
        crop = image.crop((0, 0, image.width, int(image.height * 0.54)))
        crop.save(qc_crop)

    annual = pd.read_csv(VALIDATION / "annual_yield_comparison_5km_10km.csv")
    metrics = pd.read_csv(VALIDATION / "cross_year_validation_metrics.csv")
    resolution = pd.read_csv(VALIDATION / "resolution_10km_vs_5km.csv")

    cal_asset = ASSETS / "yield_calibration_2020.png"
    d2020 = annual[(annual.resolution_m == 5000) & (annual.year == 2020)].set_index("crop")
    crops, labels = ["wheat", "maize"], ["小麦", "玉米"]
    x = np.arange(2)
    fig, ax = plt.subplots(figsize=(6.7, 3.5))
    ax.bar(x - 0.18, [d2020.loc[c, "raw_apsim_yield_kg_ha"] for c in crops], 0.36,
           color=rgb_mpl(BLUE), label="原始APSIM")
    ax.bar(x + 0.18, [d2020.loc[c, "official_yield_kg_ha"] for c in crops], 0.36,
           color=rgb_mpl(TEAL), label="正式县均")
    for i, crop in enumerate(crops):
        raw = d2020.loc[crop, "raw_apsim_yield_kg_ha"]
        obs = d2020.loc[crop, "official_yield_kg_ha"]
        ax.text(i - 0.18, raw + 250, f"{raw:,.0f}", ha="center", fontsize=8)
        ax.text(i + 0.18, obs + 250, f"{obs:,.0f}", ha="center", fontsize=8)
    ax.set_xticks(x, labels)
    ax.set_ylabel("产量 (kg ha$^{-1}$)")
    ax.set_title("2020年县域面积加权产量")
    ax.legend(frameon=False, ncol=2, loc="upper left")
    ax.set_ylim(0, 14500)
    fig.tight_layout()
    fig.savefig(cal_asset, dpi=260, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    trend_asset = ASSETS / "multiyear_yield_validation.png"
    d5 = annual[annual.resolution_m == 5000]
    fig, axes = plt.subplots(1, 2, figsize=(8.2, 3.25), sharex=True)
    for ax, crop, label in zip(axes, crops, labels):
        group = d5[d5.crop == crop].sort_values("year")
        ax.plot(group.year, group.official_yield_kg_ha, "o-", color=rgb_mpl(TEAL), lw=2, label="正式县均")
        ax.plot(group.year, group.raw_apsim_yield_kg_ha, "o--", color=rgb_mpl(BLUE), lw=1.7, label="原始APSIM")
        ax.plot(group.year, group.calibrated_yield_kg_ha, "o-.", color=rgb_mpl(GOLD), lw=1.7, label="固定2020系数")
        ax.set_title(label)
        ax.set_xticks([2018, 2019, 2020])
        ax.set_ylabel("kg ha$^{-1}$")
    axes[0].legend(frameon=False, fontsize=7.5, loc="upper left")
    fig.tight_layout()
    fig.savefig(trend_asset, dpi=260, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    mape_asset = ASSETS / "crossyear_mape.png"
    m5 = metrics[metrics.resolution_m == 5000].set_index("crop")
    fig, ax = plt.subplots(figsize=(6.5, 3.25))
    ax.bar(x - 0.18, [m5.loc[c, "raw_mape_percent"] for c in crops], 0.36,
           color=rgb_mpl(BLUE), label="原始")
    ax.bar(x + 0.18, [m5.loc[c, "calibrated_mape_percent"] for c in crops], 0.36,
           color=rgb_mpl(TEAL), label="固定系数后")
    for xpos, value in zip(list(x - 0.18) + list(x + 0.18),
                            [m5.loc[c, "raw_mape_percent"] for c in crops] + [m5.loc[c, "calibrated_mape_percent"] for c in crops]):
        ax.text(xpos, value + 1.2, f"{value:.1f}%", ha="center", fontsize=8)
    ax.set_xticks(x, labels)
    ax.set_ylabel("2018–2019 MAPE (%)")
    ax.set_ylim(0, 56)
    ax.legend(frameon=False, ncol=2)
    fig.tight_layout()
    fig.savefig(mape_asset, dpi=260, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    res_asset = ASSETS / "resolution_effect_percent.png"
    fig, ax = plt.subplots(figsize=(7.0, 3.2))
    years = [2018, 2019, 2020]
    wheat = resolution[resolution.crop == "wheat"].set_index("year").loc[years, "raw_10km_minus_5km_percent"]
    maize = resolution[resolution.crop == "maize"].set_index("year").loc[years, "raw_10km_minus_5km_percent"]
    xx = np.arange(3)
    ax.bar(xx - 0.18, wheat, 0.36, color=rgb_mpl(BLUE), label="小麦")
    ax.bar(xx + 0.18, maize, 0.36, color=rgb_mpl(TEAL), label="玉米")
    ax.axhline(0, color="#555555", lw=0.8)
    ax.set_xticks(xx, years)
    ax.set_ylabel("10 km相对5 km变化 (%)")
    ax.set_ylim(-0.48, 0.20)
    ax.legend(frameon=False, ncol=2)
    fig.tight_layout()
    fig.savefig(res_asset, dpi=260, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    return {
        "map": map_asset,
        "qc_crop": qc_crop,
        "calibration": cal_asset,
        "trend": trend_asset,
        "mape": mape_asset,
        "resolution": res_asset,
    }


def rgb_mpl(value: str) -> str:
    return f"#{value}"


def add_text(slide, text, x, y, w, h, size=16, color=INK, bold=False,
             align=PP_ALIGN.LEFT, font="Microsoft YaHei", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_bullets(slide, items, x, y, w, h, size=14, color=INK, spacing=7):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(spacing)
        p.line_spacing = 1.12
        p._p.get_or_add_pPr().insert(0, p._p._new_buChar()) if False else None
    return box


def add_rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    return shape


def add_picture_contain(slide, path: Path, x, y, w, h):
    with Image.open(path) as image:
        ratio = image.width / image.height
    box_ratio = w / h
    if ratio >= box_ratio:
        pw, ph = w, w / ratio
        px, py = x, y + (h - ph) / 2
    else:
        ph, pw = h, h * ratio
        py, px = y, x + (w - pw) / 2
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def title(slide, text, number=None):
    add_text(slide, text, 0.55, 0.28, 11.8, 0.55, size=25, color=NAVY, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.92), Inches(1.0), Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = rgb(TEAL); line.line.fill.background()
    if number is not None:
        add_text(slide, f"{number:02d}", 12.25, 0.3, 0.45, 0.35, size=9, color=GREY, align=PP_ALIGN.RIGHT)


def footer(slide, source="项目内部结果与可追溯数据文件"):
    add_text(slide, source, 0.58, 7.15, 10.8, 0.18, size=7.2, color=GREY)


def takeaway(slide, text):
    add_rect(slide, 0.55, 6.55, 12.2, 0.42, PALE_TEAL, radius=False)
    add_text(slide, text, 0.75, 6.61, 11.8, 0.25, size=11.5, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_metric(slide, value, label, x, y, w=2.7, color=TEAL):
    add_rect(slide, x, y, w, 1.25, LIGHT, line="E2E6E8")
    add_text(slide, value, x + 0.12, y + 0.14, w - 0.24, 0.48, size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.12, y + 0.73, w - 0.24, 0.28, size=10, color=GREY, align=PP_ALIGN.CENTER)


def add_table(slide, data, x, y, w, h, header_fill=NAVY, font_size=10):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text = str(data[row][col])
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.fill.solid(); cell.fill.fore_color.rgb = rgb(header_fill if row == 0 else (WHITE if row % 2 else LIGHT))
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(font_size if row else font_size + 0.5)
                p.font.bold = row == 0
                p.font.color.rgb = rgb(WHITE if row == 0 else INK)
                p.text_frame if False else None
    return table


def build_deck(assets: dict[str, Path]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, WHITE, radius=False)
    add_rect(slide, 0, 0, 4.75, 7.5, NAVY, radius=False)
    add_text(slide, "齐河县县域APSIM格网模拟", 0.62, 1.15, 3.55, 1.1, 27, WHITE, True)
    add_text(slide, "数据构建、绝对产量校准与分辨率敏感性", 0.62, 2.38, 3.55, 1.05, 18, "DCE7EE", True)
    add_text(slide, "2020轮作基线 · 2018–2020跨年验证 · 1/5/10 km格网", 0.62, 3.72, 3.55, 0.72, 11.5, "BFD3DF")
    add_text(slide, "近期县域工作总结", 0.62, 6.55, 3.2, 0.35, 10.5, WHITE)
    add_picture_contain(slide, assets["map"], 4.98, 0.68, 8.0, 5.9)
    add_text(slide, "齐河县2020年小麦—玉米轮作区", 6.3, 6.6, 5.1, 0.3, 9, GREY, align=PP_ALIGN.CENTER)
    add_notes(slide, "本次汇报总结近期齐河县县域尺度APSIM工作。主线是从空间数据闭环出发，建立可复现的轮作模拟单元，完成绝对产量校准、跨年验证和格网分辨率敏感性实验。")

    # 2 executive summary
    slide = prs.slides.add_slide(blank); title(slide, "已经形成可运行、可校准、可比较的县域模拟链", 2)
    add_metric(slide, "58,918.84 ha", "遥感轮作制图面积", 0.75, 1.45, 2.75, TEAL)
    add_metric(slide, "101", "5 km唯一APSIM组合", 3.78, 1.45, 2.75, BLUE)
    add_metric(slide, "−0.39%", "2020小麦原始偏差", 6.81, 1.45, 2.75, TEAL)
    add_metric(slide, "+65.54%", "2020玉米原始偏差", 9.84, 1.45, 2.75, RED)
    add_bullets(slide, [
        "气象升级为AgERA5 v2（0.1°，2017–2020），土壤按格网内HWSD类型面积拆分。",
        "2020正式统计用于建立固定系数；2018–2019作为跨年检验，避免同年自我验证。",
        "10 km相对5 km的县均产量最大变化仅0.388%，但空间异质性表达有所减弱。",
    ], 1.05, 3.35, 11.1, 2.25, size=15, spacing=12)
    takeaway(slide, "当前成果已经从“能运行”推进到“有统计约束、可审计尺度效应”。")
    footer(slide); add_notes(slide, "四个数字概括阶段成果。小麦已经接近县域统计，玉米原始结果明显更接近高产示范水平，因此校准重点集中在玉米。")

    # 3 problem/gap
    slide = prs.slides.add_slide(blank); title(slide, "县域模拟的瓶颈不是单次APSIM，而是空间输入与统计口径闭环", 3)
    columns = [
        ("空间代表性", ["气象节点粗于格网", "格网内包含多种土壤", "轮作边界来自遥感分类"], BLUE),
        ("管理代表性", ["普通农户与高产示范差异", "播期、密度、N和灌溉需情景化", "品种参数可能偏向试验田"], TEAL),
        ("验证口径", ["年鉴出版年≠统计年份", "县均产量≠示范田测产", "总产量比较必须统一播种面积"], GOLD),
    ]
    for i, (head, items, color) in enumerate(columns):
        x = 0.75 + i * 4.15
        add_rect(slide, x, 1.45, 3.75, 4.45, LIGHT, line="E1E5E8")
        add_rect(slide, x, 1.45, 3.75, 0.58, color, radius=False)
        add_text(slide, head, x + 0.18, 1.57, 3.35, 0.3, 15, WHITE, True, PP_ALIGN.CENTER)
        add_bullets(slide, items, x + 0.32, 2.35, 3.1, 2.85, 13.5, spacing=14)
    takeaway(slide, "解决策略：同年数据闭合 + 土壤面积拆分 + 管理情景化 + 正式统计校准。")
    footer(slide); add_notes(slide, "这里强调研究缺口。县域工作不能只展示一张产量图，还要说明每个空间单元如何被定义、如何加权，以及验证数据到底代表县均还是高产样点。")

    # 4 data chain
    slide = prs.slides.add_slide(blank); title(slide, "数据链将轮作区、土壤、气象和管理映射到同一模拟单元", 4)
    labels = [
        ("轮作范围", "ChinaCP-Wheat10m\n2020 · 10 m", BLUE),
        ("土壤", "HWSD v2\n格网内面积拆分", TEAL),
        ("气象", "AgERA5 v2\n2017–2020", GOLD),
        ("管理", "普通/规范/示范\n播期·密度·N·灌溉", RED),
    ]
    for i, (head, body, color) in enumerate(labels):
        x = 0.65 + i * 3.05
        add_rect(slide, x, 1.4, 2.55, 1.35, LIGHT, line=color)
        add_text(slide, head, x + 0.12, 1.58, 2.3, 0.3, 14, color, True, PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.12, 2.02, 2.3, 0.52, 10.5, INK, align=PP_ALIGN.CENTER)
    add_text(slide, "↓  空间叠加与唯一组合去重", 4.6, 3.02, 4.2, 0.38, 15, NAVY, True, PP_ALIGN.CENTER)
    add_rect(slide, 2.15, 3.65, 9.05, 1.12, PALE_BLUE, line=BLUE)
    add_text(slide, "格网 × HWSD土壤子单元 × 最近AgERA5节点 × 管理情景", 2.45, 3.96, 8.45, 0.42, 18, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, "↓  APSIM Classic连续轮作：冬小麦 → 夏玉米", 3.45, 5.02, 6.5, 0.38, 14, TEAL, True, PP_ALIGN.CENTER)
    takeaway(slide, "模拟次数由唯一“土壤—气象”组合决定，结果再映射回所有代表格网。")
    footer(slide, "数据：ChinaCP-Wheat10m、HWSD v2、AgERA5 v2；管理参数来自整理后的情景配置")
    add_notes(slide, "数据链中面积不是APSIM的生长参数。APSIM计算每公顷产量，空间面积在模拟后用于生产量和县域平均的汇总。")

    # 5 rotation logic
    slide = prs.slides.add_slide(blank); title(slide, "2020轮作区直接采用公开轮作类别，不由异源单作物图强制求交", 5)
    add_picture_contain(slide, assets["qc_crop"], 0.55, 1.18, 8.35, 4.85)
    add_rect(slide, 9.15, 1.35, 3.55, 4.55, LIGHT, line="E1E5E8")
    add_text(slide, "当前判定", 9.45, 1.62, 2.95, 0.35, 15, NAVY, True)
    add_text(slide, "县界内 ∩ 有效像元 ∩\nChinaCP轮作类别 = 1", 9.45, 2.15, 2.9, 0.78, 15, TEAL, True)
    add_text(slide, "独立冬小麦/玉米图的作用", 9.45, 3.25, 2.95, 0.35, 13, NAVY, True)
    add_bullets(slide, ["边际播种面积统计", "与年鉴及轮作图交叉检查", "未来单作物情景构建"], 9.45, 3.72, 2.85, 1.48, 11.5, spacing=6)
    takeaway(slide, "不同分类体系的不一致被保留为质量诊断，不通过裁剪人为改成一致。")
    footer(slide, "来源：ChinaCP-Wheat10m 2020；CN-Wheat10；CCD-Maize；CACD-v1")
    add_notes(slide, "需要特别区分2020与2024流程。2020有直接发布的轮作类别，因此不再用独立小麦和玉米图简单相交；独立产品只用于边际面积和一致性诊断。")

    # 6 weighting
    slide = prs.slides.add_slide(blank); title(slide, "面积比例不改变单次生长过程，但决定格网和县域汇总权重", 6)
    add_rect(slide, 0.75, 1.35, 5.65, 4.65, LIGHT, line="DFE5E8")
    add_text(slide, "APSIM内部", 1.05, 1.65, 2.0, 0.35, 15, NAVY, True)
    add_bullets(slide, ["输入：土壤、气象、管理、品种", "输出：单位面积产量 (kg/ha)", "不知道代表10 ha还是1,000 ha"], 1.05, 2.18, 4.65, 1.7, 13, spacing=10)
    add_text(slide, "面积不进入作物生长方程", 1.05, 4.65, 4.9, 0.45, 18, RED, True, PP_ALIGN.CENTER)
    add_rect(slide, 6.9, 1.35, 5.65, 4.65, PALE_TEAL, line=TEAL)
    add_text(slide, "模拟后汇总", 7.2, 1.65, 2.2, 0.35, 15, NAVY, True)
    add_text(slide, "子单元生产量 = 产量 × 轮作面积", 7.2, 2.35, 4.9, 0.45, 15, TEAL, True, PP_ALIGN.CENTER)
    add_text(slide, "县均产量 = Σ(产量 × 面积) / Σ面积", 7.2, 3.15, 4.9, 0.45, 15, TEAL, True, PP_ALIGN.CENTER)
    add_text(slide, "格网内不同HWSD土壤分别模拟后，按 soil_rotation_area_ha 加权。", 7.35, 4.25, 4.55, 0.78, 12.5, INK, align=PP_ALIGN.CENTER)
    takeaway(slide, "rotation_fraction用于得到代表面积；真正进入汇总的是轮作面积权重。")
    footer(slide); add_notes(slide, "举例：同一格网中两种土壤分别代表240和160公顷，APSIM各自产生单位面积产量，最后按240比160加权。")

    # 7 baseline design
    slide = prs.slides.add_slide(blank); title(slide, "修正后的5 km普通农户基线保留土壤异质性与连续轮作过程", 7)
    add_metric(slide, "85", "有轮作面积的5 km格网", 0.72, 1.28, 2.35, BLUE)
    add_metric(slide, "248", "格网—土壤子单元", 3.28, 1.28, 2.35, TEAL)
    add_metric(slide, "14", "HWSD制图单元", 5.84, 1.28, 2.35, GOLD)
    add_metric(slide, "27", "实际使用AgERA5节点", 8.40, 1.28, 2.35, BLUE)
    add_metric(slide, "101", "唯一APSIM组合", 10.96, 1.28, 1.65, TEAL)
    table_data = [
        ["参数", "冬小麦", "夏玉米"],
        ["播期", "10月5–15日", "6月12–20日"],
        ["密度", "210万株/ha", "60,000株/ha"],
        ["N投入", "180 kg/ha", "180 kg/ha"],
        ["灌溉", "2次，共150 mm", "1次，60 mm"],
        ["收获", "模型成熟/收获", "10月1日"],
    ]
    add_table(slide, table_data, 1.05, 3.15, 7.4, 2.65, font_size=10.5)
    add_rect(slide, 8.85, 3.15, 3.45, 2.65, LIGHT, line="E0E5E7")
    add_text(slide, "气象窗口", 9.15, 3.47, 2.85, 0.3, 13, NAVY, True)
    add_text(slide, "2017–2020完整日值\n目标轮作季：2019-10至2020-12", 9.15, 3.94, 2.8, 0.8, 12, INK)
    add_text(slide, "总灌溉：210 mm", 9.15, 5.05, 2.8, 0.35, 14, TEAL, True)
    takeaway(slide, "不先平均土壤参数，避免水氮过程的非线性被空间均值掩盖。")
    footer(slide); add_notes(slide, "5 km只是外层格网。格网内部按轮作像元对应的HWSD类型拆为多个子单元，所以实际模拟单元多于格网数。")

    # 8 calibration
    slide = prs.slides.add_slide(blank); title(slide, "2020小麦已接近正式县均，玉米原始结果明显偏向高产水平", 8)
    add_picture_contain(slide, assets["calibration"], 0.65, 1.25, 7.6, 4.65)
    add_rect(slide, 8.55, 1.35, 4.15, 4.45, LIGHT, line="E1E5E8")
    add_text(slide, "固定统计校准系数", 8.88, 1.72, 3.45, 0.35, 15, NAVY, True)
    add_text(slide, "小麦  1.003874", 9.05, 2.45, 3.1, 0.45, 20, TEAL, True, PP_ALIGN.CENTER)
    add_text(slide, "玉米  0.604069", 9.05, 3.18, 3.1, 0.45, 20, RED, True, PP_ALIGN.CENTER)
    add_bullets(slide, ["不修改已有品种遗传参数", "保留原始APSIM与校准结果两套输出", "同年校准后零残差不等于独立验证"], 8.95, 4.08, 3.05, 1.35, 10.8, spacing=4)
    takeaway(slide, "玉米高估更可能反映品种/管理接近高产示范，而非县域普通生产平均。")
    footer(slide, "正式统计：《德州统计年鉴2021》，齐河县2020年作物生产表")
    add_notes(slide, "小麦原始偏差只有负0.39%，没有必要为了一个县年数据重拟品种。玉米高出65.54%，采用统计后校准而不直接破坏已有高产点位参数。")

    # 9 multiyear
    slide = prs.slides.add_slide(blank); title(slide, "2018–2020正式统计建立了跨年检验框架", 9)
    add_picture_contain(slide, assets["trend"], 0.55, 1.2, 8.55, 4.8)
    stats = pd.read_csv(STATS)
    stats_table = [["年份", "小麦", "玉米"]]
    for year in (2018, 2019, 2020):
        row = stats[stats.year == year].set_index("crop")
        stats_table.append([year, f"{row.loc['wheat','yield_kg_ha']:,.0f}", f"{row.loc['maize','yield_kg_ha']:,.0f}"])
    add_table(slide, stats_table, 9.35, 1.6, 3.05, 2.15, font_size=10)
    add_text(slide, "单位：kg/ha", 9.42, 3.85, 2.9, 0.22, 8, GREY, align=PP_ALIGN.RIGHT)
    add_rect(slide, 9.35, 4.25, 3.05, 1.45, PALE_BLUE, line=BLUE)
    add_text(slide, "验证设计", 9.62, 4.48, 2.5, 0.3, 13, NAVY, True)
    add_text(slide, "2020拟合系数\n2018–2019固定系数外推", 9.62, 4.92, 2.5, 0.58, 11, INK, align=PP_ALIGN.CENTER)
    takeaway(slide, "跨年检验关注系数可迁移性，而不是逐年重新调到零误差。")
    footer(slide, "正式统计：《德州统计年鉴2019–2021》；年鉴年份分别对应上一统计年度")
    add_notes(slide, "2019和2020年鉴截图最初容易被误当成同年数据。这里已经按年鉴出版年减一的统计年度重新核准。")

    # 10 validation
    slide = prs.slides.add_slide(blank); title(slide, "固定2020系数显著改善玉米跨年误差，但不能完全吸收年际变化", 10)
    add_picture_contain(slide, assets["mape"], 0.65, 1.3, 7.25, 4.55)
    add_rect(slide, 8.25, 1.45, 4.2, 4.25, LIGHT, line="E1E5E8")
    add_text(slide, "2018–2019校准后偏差", 8.62, 1.82, 3.45, 0.35, 15, NAVY, True)
    add_table(slide, [
        ["年份", "小麦", "玉米"],
        ["2018", "−7.48%", "−6.62%"],
        ["2019", "−6.32%", "−14.34%"],
    ], 8.65, 2.45, 3.35, 1.6, font_size=10.5)
    add_text(slide, "三年联合诊断系数", 8.65, 4.38, 3.35, 0.28, 12, NAVY, True)
    add_text(slide, "小麦 ≈ 1.0495   玉米 ≈ 0.6467", 8.65, 4.83, 3.35, 0.38, 14, TEAL, True, PP_ALIGN.CENTER)
    takeaway(slide, "玉米MAPE由48.20%降至10.48%；三年联合系数仅作诊断，不能再称独立验证。")
    footer(slide); add_notes(slide, "玉米固定系数方向正确且改进明显，但2019仍低估14.34%，说明真实管理、品种或统计含水量等因素具有年际变化。")

    # 11 maps
    slide = prs.slides.add_slide(blank); title(slide, "1、5、10 km聚合保持轮作总面积一致，但空间细节逐级减少", 11)
    add_picture_contain(slide, assets["map"], 0.45, 1.08, 12.45, 5.65)
    footer(slide, "颜色表示格网内轮作面积比例；三个尺度轮作总面积均为58,918.84 ha")
    add_notes(slide, "1 km包含1558个县域格网，5 km为85个，10 km为28个。颜色是轮作面积除以格网与县界相交面积，不是APSIM产量。")

    # 12 resolution metrics
    slide = prs.slides.add_slide(blank); title(slide, "县域均值对5→10 km不敏感，计算量仅小幅下降", 12)
    add_picture_contain(slide, assets["resolution"], 0.55, 1.25, 7.25, 4.6)
    add_table(slide, [
        ["指标", "5 km", "10 km"],
        ["县域格网", "85", "28"],
        ["土壤子单元", "248", "111"],
        ["唯一APSIM组合", "101", "88"],
        ["运行时间", "452.43 s", "390.66 s"],
    ], 8.15, 1.55, 4.05, 2.75, font_size=10.5)
    add_rect(slide, 8.15, 4.62, 4.05, 1.05, PALE_TEAL, line=TEAL)
    add_text(slide, "最大县均变化", 8.42, 4.83, 1.8, 0.25, 11, GREY, True)
    add_text(slide, "0.388%", 10.18, 4.72, 1.55, 0.45, 24, TEAL, True, PP_ALIGN.CENTER)
    takeaway(slide, "10 km适合县域总量与快速情景；5 km更适合表达县内空间异质性。")
    footer(slide); add_notes(slide, "气象原始分辨率为0.1度，约9到11公里，这也是5和10公里县均差异很小的重要原因。10公里减少了格网和土壤子单元，但唯一土壤气象组合只减少约13%。")

    # 13 limitations/next
    slide = prs.slides.add_slide(blank); title(slide, "下一阶段应从“固定空间+统计校准”转向多年份真实空间验证", 13)
    cards = [
        ("当前边界", ["2018–2019仍使用2020轮作掩膜", "管理情景跨年保持不变", "县域统计不能验证像元位置"], RED),
        ("优先补充", ["逐年轮作/单作物空间边界", "乡镇或田间产量与物候", "统计含水量和收获损失口径"], GOLD),
        ("推荐实验", ["1/5/10 km统一固定系数", "多年份留一法或独立测试年", "气象、土壤、管理贡献分解"], TEAL),
    ]
    for i, (head, items, color) in enumerate(cards):
        x = 0.72 + i * 4.18
        add_rect(slide, x, 1.42, 3.78, 4.35, LIGHT, line="E0E5E7")
        add_rect(slide, x, 1.42, 3.78, 0.6, color, radius=False)
        add_text(slide, head, x + 0.15, 1.57, 3.48, 0.28, 14, WHITE, True, PP_ALIGN.CENTER)
        add_bullets(slide, items, x + 0.34, 2.35, 3.1, 2.45, 12.5, spacing=11)
    takeaway(slide, "阶段结论：空间流程已经稳定；主要科学不确定性转向年度管理、品种代表性与本地观测。")
    footer(slide); add_notes(slide, "最后强调下一步不是继续增加格网数量，而是增加独立年份和本地观测。空间分辨率实验已经说明县域均值对5到10公里并不敏感。")

    return prs


def validate(prs: Presentation, assets: dict[str, Path]) -> dict:
    prs.save(OUT)
    reopened = Presentation(OUT)
    media = list(__import__("zipfile").ZipFile(OUT).namelist())
    media_files = [item for item in media if item.startswith("ppt/media/")]
    notes_count = 0
    for slide in reopened.slides:
        try:
            if slide.notes_slide.notes_text_frame.text.strip():
                notes_count += 1
        except Exception:
            pass
    return {
        "pptx": str(OUT),
        "slides": len(reopened.slides),
        "embedded_media_files": len(media_files),
        "slides_with_notes": notes_count,
        "assets": {name: str(path) for name, path in assets.items()},
        "reopen_validation": "PASS",
        "package_media_validation": "PASS" if len(media_files) >= len(assets) else "CHECK",
    }


def write_supporting_files(result: dict, assets: dict[str, Path]) -> None:
    QA.write_text(
        "# PPTX QA report\n\n"
        f"- PPTX creation: PASS\n- Slide count: {result['slides']}\n"
        f"- Embedded media files: {result['embedded_media_files']}\n"
        f"- Slides with speaker notes: {result['slides_with_notes']}\n"
        "- Verification: reopened with python-pptx; ZIP package and embedded media inspected.\n"
        "- Placeholders: none.\n"
        "- Known limitation: no LibreOffice/PowerPoint headless renderer was available, so QA is structural rather than rendered slide-by-slide.\n"
        "- Manual follow-up: open in PowerPoint once to confirm local font substitution and animation-free layout.\n",
        encoding="utf-8",
    )
    rows = ["# Asset manifest", "", "| Asset | Source | Method | Slides |", "|---|---|---|---|"]
    mapping = {
        "map": ("Existing 1/5/10 km grid figure", "copied without data alteration", "1, 11"),
        "qc_crop": ("Existing crop/grid QC figure", "top-panel crop only", "5"),
        "calibration": ("2020 annual yield comparison CSV", "matplotlib redraw", "8"),
        "trend": ("2018–2020 annual yield comparison CSV", "matplotlib redraw", "9"),
        "mape": ("cross-year validation metrics CSV", "matplotlib redraw", "10"),
        "resolution": ("10 km versus 5 km comparison CSV", "matplotlib redraw", "12"),
    }
    for name, path in assets.items():
        source, method, slides = mapping[name]
        rows.append(f"| `{path.name}` | {source} | {method} | {slides} |")
    MANIFEST.write_text("\n".join(rows) + "\n", encoding="utf-8")


def main() -> None:
    REPORT.mkdir(parents=True, exist_ok=True)
    assets = setup_assets()
    prs = build_deck(assets)
    result = validate(prs, assets)
    write_supporting_files(result, assets)
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
