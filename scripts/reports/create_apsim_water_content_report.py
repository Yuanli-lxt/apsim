from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parents[2] / "outputs" / "reports" / "apsim_water_content_report_after_0510.pptx"
FONT = "Microsoft YaHei"
BG = RGBColor(248, 249, 247)
INK = RGBColor(35, 40, 45)
MUTED = RGBColor(102, 112, 120)
GREEN = RGBColor(42, 120, 92)
BLUE = RGBColor(40, 88, 145)
ORANGE = RGBColor(190, 105, 40)
RED = RGBColor(170, 70, 70)
LINE = RGBColor(210, 216, 214)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, x, y, w, h, text="", size=20, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.75, 0.42, 11.8, 0.5, title, 25, True, INK)
    if subtitle:
        add_textbox(slide, 0.77, 0.93, 11.5, 0.35, subtitle, 10.5, False, MUTED)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.22), Inches(11.9), Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.color.rgb = LINE


def bullet_block(slide, x, y, w, h, bullets, size=15, color=INK, gap=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.space_after = Pt(8 if gap else 4)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def metric_card(slide, x, y, w, h, value, label, color=GREEN):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shp.line.color.rgb = LINE
    add_textbox(slide, x + 0.18, y + 0.15, w - 0.36, 0.42, value, 24, True, color, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.15, y + 0.68, w - 0.3, 0.38, label, 10.5, False, MUTED, PP_ALIGN.CENTER)


def add_table(slide, x, y, w, h, data, col_widths=None, font_size=10.5):
    table_shape = slide.shapes.add_table(len(data), len(data[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)
    for row_idx, row in enumerate(data):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(229, 236, 232) if row_idx == 0 else RGBColor(255, 255, 255)
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(font_size)
                p.font.bold = row_idx == 0
                p.font.color.rgb = INK
    return table_shape


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(slide, 0.78, 1.05, 11.5, 0.6, "APSIM 土壤含水量预测改进汇报", 34, True, INK)
    add_textbox(slide, 0.82, 1.85, 10.5, 0.45, "5 月 10 日之后的模型、搜索与验证变化", 18, False, MUTED)
    metric_card(slide, 0.82, 3.05, 2.35, 1.25, "0.2535", "当前 soil_water error", GREEN)
    metric_card(slide, 3.45, 3.05, 2.35, 1.25, "0.548", "InitialWater.FractionFull", BLUE)
    metric_card(slide, 6.08, 3.05, 2.35, 1.25, "0.56", "crit_fr_asw", ORANGE)
    metric_card(slide, 8.71, 3.05, 2.35, 1.25, "6 天", "物候最大误差", GREEN)
    add_textbox(
        slide,
        0.85,
        6.65,
        8.5,
        0.3,
        "数据来源：process_bio git 记录、output_sobol / output_hdsw / output_hdsw_sobol_water_yield 结果文件",
        9.5,
        False,
        MUTED,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "一句话结论")
    bullet_block(
        slide,
        0.9,
        1.65,
        11.2,
        3.7,
        [
            "5.10 后，工作重心从“能跑 + 生物量校准”转向“土壤含水量预测误差控制”。",
            "HDSW 土壤替换路线暴露出水分库与产量过程不匹配，短期不适合作为主线。",
            "当前主线为 Sobol 品种校准 + water-yield 约束搜索 + InitialWater 局部精细搜索。",
            "候选接受规则从单一综合分，升级为 soil_water 优先、yield 硬约束、phenology 守门。",
        ],
        18,
        INK,
        True,
    )
    metric_card(slide, 0.95, 5.65, 3.2, 1.05, "2110 条", "土壤含水量观测参与评分", BLUE)
    metric_card(slide, 4.55, 5.65, 3.2, 1.05, "0 缺失", "当前 best 模拟值缺失", GREEN)
    metric_card(slide, 8.15, 5.65, 3.2, 1.05, "10-50 cm", "分层 soil_water 诊断", ORANGE)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "5.10 后时间线")
    add_table(
        slide,
        0.75,
        1.55,
        11.85,
        4.65,
        [
            ["日期", "阶段", "主要变化"],
            ["5.13", "HWSD/HDSW 接入", "土壤转换、预测-观测对比、HDSW no-soil-edit 搜索"],
            ["5.17", "首次提交", "基础脚本、Sobol 工作流、评估脚本固化"],
            ["5.18", "Sobol 后迭代", "扩展输出、Sobol 指数、自动汇报脚本、主搜索增强"],
            ["5.19", "水分-产量约束", "water_yield_search、HDSW 诊断、InitialWater/crit 局部搜索"],
            ["5.25", "sobol_apsim", "固化水分-产量搜索和本地验证结果"],
        ],
        [1.05, 2.05, 8.75],
        11.3,
    )
    add_textbox(slide, 0.85, 6.45, 11.2, 0.35, "汇报口径：从“可复现实验流程”逐步升级到“含水量预测优化流程”。", 12.5, True, GREEN)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "评估对象变了")
    bullet_block(
        slide,
        0.85,
        1.55,
        5.75,
        4.3,
        [
            "原来：更偏向小麦/玉米生物量、产量和物候的综合校准。",
            "现在：soil_water 成为独立关键目标，并进入每轮候选接受判断。",
            "每轮输出 prediction_vs_truth.csv、summary、metrics 和 stage_alignment。",
            "LAI 保留为诊断项，核心权重集中在生物量、结构、产量、soil_water 和物候。",
        ],
        14.5,
        INK,
        True,
    )
    add_table(
        slide,
        7.05,
        1.55,
        4.95,
        3.45,
        [
            ["评分组", "当前 best mean rel error"],
            ["total_biomass", "0.2572"],
            ["structure", "0.4239"],
            ["yield", "0.1266"],
            ["soil_water", "0.2535"],
            ["phenology", "0.1500"],
        ],
        [2.55, 2.4],
        12,
    )
    add_textbox(slide, 7.1, 5.35, 4.8, 0.65, "soil_water 对应 10-50 cm 分层观测，共 2110 条，当前 best 无模拟缺失。", 12, False, MUTED)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Sobol 进入主流程")
    bullet_block(
        slide,
        0.85,
        1.45,
        11.35,
        4.5,
        [
            "清点 cultivar 参数 -> 构建参数范围 -> 生成 Sobol/Saltelli 样本。",
            "批量生成 APSIM runs -> 运行 APSIM -> 汇总输出 -> 计算 Sobol 指数。",
            "在主搜索中用敏感性 + signed error 选择阶段和参数。",
            "典型例子：maize_phenology 阶段调低 tt_flower_to_maturity，物候最大误差控制到 6 天。",
        ],
        17,
        INK,
        True,
    )
    metric_card(slide, 1.0, 5.75, 3.2, 1.0, "0.2359", "Sobol iter 3 综合分", GREEN)
    metric_card(slide, 4.6, 5.75, 3.2, 1.0, "0.2795", "Sobol iter 3 soil_water", ORANGE)
    metric_card(slide, 8.2, 5.75, 3.2, 1.0, "0.0893", "Sobol iter 3 yield", BLUE)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Water-yield 搜索：把含水量作为主目标")
    bullet_block(
        slide,
        0.85,
        1.5,
        6.0,
        4.65,
        [
            "主目标：降低 soil_water 分层误差。",
            "硬约束：wheat / maize yield error 必须在阈值内。",
            "守门条件：phenology 不能突破阈值。",
            "稳定性惩罚：防止水分改善以生物量结构崩坏为代价。",
            "诊断显示 10-50 cm 多层偏湿，因此降低灌溉阈值。",
        ],
        14.5,
        INK,
        True,
    )
    add_table(
        slide,
        7.15,
        1.7,
        4.8,
        3.25,
        [
            ["动作", "变化"],
            ["crit_fr_asw", "0.60 -> 0.56"],
            ["soil_water error", "0.2795 -> 0.2625"],
            ["water_yield_score", "0.2397 -> 0.2373"],
            ["yield/phenology", "均通过守门"],
        ],
        [2.35, 2.45],
        12.5,
    )
    add_textbox(slide, 7.2, 5.3, 4.65, 0.7, "这一步的意义：不追求综合分瞬时最优，而是让含水量预测朝正确方向移动。", 12.5, True, GREEN)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "HDSW 路线的诊断结论")
    metric_card(slide, 0.9, 1.55, 2.55, 1.15, "1.4295", "HDSW soil_water error", RED)
    metric_card(slide, 3.75, 1.55, 2.55, 1.15, "0.7474", "wheat yield error", RED)
    metric_card(slide, 6.6, 1.55, 2.55, 1.15, "1.0000", "maize yield error", RED)
    metric_card(slide, 9.45, 1.55, 2.55, 1.15, "0", "maize yield 瓶颈", RED)
    bullet_block(
        slide,
        0.9,
        3.35,
        11.2,
        2.6,
        [
            "HDSW soil 基线下，maize yield = 0，wheat yield 也明显低于观测。",
            "InitialWater、crit_fr_asw、灌溉模式和小步品种扰动均未能在不恶化其他目标的情况下恢复产量。",
            "汇报口径：HDSW 不是“没用”，而是证明当前土壤物理库直接替换会破坏作物-水分-产量耦合。",
        ],
        15.5,
        INK,
        True,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "局部精细搜索确定当前 best")
    add_table(
        slide,
        0.85,
        1.45,
        5.65,
        4.5,
        [
            ["项目", "结果"],
            ["FractionFull", "0.554 -> 0.548"],
            ["crit_fr_asw", "0.56 固定"],
            ["maize tt_flag_to_flower", "38.3 -> 38.683"],
            ["soil_water error", "0.2547 -> 0.2535"],
            ["wheat / maize yield error", "0.1087 / 0.1445"],
            ["phenology max error", "6 天"],
        ],
        [2.55, 3.1],
        12,
    )
    bullet_block(
        slide,
        7.05,
        1.6,
        5.1,
        4.0,
        [
            "细网格测试 FractionFull = 0.546-0.552。",
            "7 个候选 APSIM 全部通过。",
            "0.546 和 0.547 明显变差，0.549-0.552 也未优于 0.548。",
            "结论：0.548 是当前局部最优，不建议继续单独压 FractionFull。",
        ],
        14.5,
        INK,
        True,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "关键结果对比")
    add_table(
        slide,
        0.65,
        1.45,
        12.1,
        4.55,
        [
            ["阶段", "soil_water", "yield", "说明"],
            ["HDSW best", "1.4348", "0.8727", "水分和产量均失配"],
            ["Sobol iter 3", "0.2795", "0.0893", "综合分较好，但含水量偏湿"],
            ["water-yield iter 4", "0.2625", "0.1207", "降低 crit_fr_asw 后水分改善"],
            ["搜索停止前 best seen", "0.2547", "约 0.13-0.15", "80 轮后无足够改善"],
            ["two-stage local best", "0.2535", "0.1087 / 0.1445", "当前可接受 best"],
        ],
        [2.45, 1.35, 1.85, 6.45],
        10.6,
    )
    add_textbox(slide, 0.85, 6.25, 11.4, 0.45, "注意：water-yield 阶段的目标是含水量预测改善，并用产量/物候守门，不等同于全真值综合分单项最优。", 11.5, True, ORANGE)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "结论与下一步")
    bullet_block(
        slide,
        0.9,
        1.5,
        5.65,
        4.5,
        [
            "当前框架已从“生物量优先搜索”升级为“含水量预测核心 + 产量/物候守门”。",
            "HDSW 土壤替换短期不作为主线，避免产量过程崩塌。",
            "FractionFull = 0.548 附近已完成细网格验证，继续单独调整收益很小。",
        ],
        15.5,
        INK,
        True,
    )
    bullet_block(
        slide,
        7.05,
        1.5,
        5.25,
        4.5,
        [
            "下一步 1：围绕 maize tt_flag_to_flower = 38.683 做更小步长安全边际测试。",
            "下一步 2：画 10-50 cm 分层 soil_water 误差图，解释哪层仍偏湿/偏干。",
            "下一步 3：若继续降低 soil_water，需要引入更明确的水分过程参数或土壤水力参数约束。",
        ],
        15.5,
        INK,
        True,
    )

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
