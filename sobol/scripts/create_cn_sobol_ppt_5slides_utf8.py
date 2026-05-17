# -*- coding: utf-8 -*-
"""Create a 5-slide UTF-8 Chinese PPTX report with verified Chinese XML."""

from __future__ import annotations

import json
import re
import zipfile
from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


BASE = Path(r"F:\APSIM710-r4221\process_bio\sobol")
N128 = BASE / "organized_outputs_screened_N128_20260515_185604" / "final_results"
OUT = BASE / "ppt_report_sobol_20260517_5slides_utf8"
ASSETS = OUT / "assets" / "figures"
FONT = "Microsoft YaHei"


def u(text: str) -> str:
    return text.encode("ascii").decode("unicode_escape")


def read_csv(path: Path) -> pd.DataFrame:
    try:
        return pd.read_csv(path, encoding="utf-8-sig")
    except UnicodeDecodeError:
        return pd.read_csv(path, encoding="gb18030")


C = {
    "navy": RGBColor(28, 51, 74),
    "blue": RGBColor(61, 104, 153),
    "teal": RGBColor(70, 145, 141),
    "green": RGBColor(78, 132, 83),
    "orange": RGBColor(196, 116, 65),
    "red": RGBColor(180, 70, 72),
    "gray": RGBColor(88, 96, 104),
    "line": RGBColor(210, 218, 226),
}


def set_font(run, size=18, bold=False, color=None):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        elem = rpr.find(qn(tag))
        if elem is None:
            elem = OxmlElement(tag)
            rpr.append(elem)
        elem.set("typeface", FONT)


def textbox(slide, x, y, w, h, text, size=18, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color or C["navy"])
    return box


def bullets(slide, x, y, w, h, items, size=19, color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(7)
        for run in p.runs:
            set_font(run, size=size, color=color or C["navy"])
    return box


def title(slide, text, subtitle=None):
    textbox(slide, 0.45, 0.25, 12.4, 0.62, text, size=30, bold=True, color=C["navy"])
    if subtitle:
        textbox(slide, 0.48, 0.88, 12.2, 0.34, subtitle, size=13, color=C["gray"])


def metric(slide, x, y, w, h, value, label, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 251)
    shape.line.color.rgb = color
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    set_font(r, size=25, bold=True, color=color)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    set_font(r2, size=12, color=C["gray"])


def takeaway(slide, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.77), Inches(12.25), Inches(0.44))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(235, 244, 242)
    shape.line.color.rgb = RGBColor(194, 218, 214)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.15)
    r = tf.paragraphs[0].add_run()
    r.text = text
    set_font(r, size=15, bold=True, color=C["navy"])


def panel(slide, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 251)
    shape.line.color.rgb = C["line"]
    return shape


def picture(slide, path, x, y, w=None, h=None):
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if w is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    if h is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def make_figures():
    ASSETS.mkdir(parents=True, exist_ok=True)
    idx = read_csv(N128 / "sobol_indices_summary.csv")
    for col in ["S1", "S1_conf", "ST", "ST_conf"]:
        idx[col] = pd.to_numeric(idx[col], errors="coerce")
    idx = idx[idx["ST"].notna()].copy()
    idx["parameter_name"] = idx["parameter_key"].astype(str).str.split("__").str[2]
    idx["parameter_name"] = idx["parameter_name"].str.replace("_1", "[1]", regex=False).str.replace("_2", "[2]", regex=False)
    stability = read_csv(N128 / "sobol_N64_vs_N128_stability_report.csv")
    stability = stability[["crop", "target_variable", "spearman_ST_rank"]].drop_duplicates()

    sns.set_theme(style="whitegrid", context="talk")
    plt.rcParams["font.family"] = "DejaVu Sans"
    palette = {"maize": "#4C78A8", "wheat": "#72B7B2"}

    heat = idx.pivot_table(index="parameter_name", columns=["crop", "target_variable"], values="ST", aggfunc="mean").fillna(0)
    heat = heat.loc[heat.max(axis=1).sort_values(ascending=False).index]
    fig, ax = plt.subplots(figsize=(8.2, 4.8))
    sns.heatmap(heat, cmap="YlGnBu", linewidths=0.25, linecolor="white", cbar_kws={"label": "ST"}, ax=ax)
    ax.set_xlabel("crop / target")
    ax.set_ylabel("parameter")
    ax.tick_params(axis="x", labelrotation=45, labelsize=8)
    ax.tick_params(axis="y", labelsize=8)
    fig.tight_layout()
    heat_path = ASSETS / "st_heatmap_en.png"
    fig.savefig(heat_path, dpi=320, bbox_inches="tight")
    plt.close(fig)

    fig, axes = plt.subplots(1, 2, figsize=(9.6, 3.6))
    for ax, crop, label in zip(axes, ["maize", "wheat"], ["maize grain_yield", "wheat grain_yield"]):
        data = idx[(idx["crop"] == crop) & (idx["target_variable"] == "grain_yield")].sort_values("ST", ascending=False).head(5)
        ax.barh(data["parameter_name"][::-1], data["ST"][::-1], color=palette[crop])
        ax.set_title(label, fontsize=12)
        ax.set_xlabel("ST")
        ax.set_ylabel("")
        ax.set_xlim(0, max(0.95 if crop == "maize" else 0.38, data["ST"].max() * 1.2))
        for i, value in enumerate(data["ST"][::-1]):
            ax.text(value + 0.01, i, f"{value:.2f}", va="center", fontsize=8)
    fig.tight_layout()
    gy_path = ASSETS / "grain_yield_top_st_combined_en.png"
    fig.savefig(gy_path, dpi=320, bbox_inches="tight")
    plt.close(fig)

    stability["group"] = stability["crop"] + "\n" + stability["target_variable"]
    fig, ax = plt.subplots(figsize=(8.0, 3.2))
    ax.bar(stability["group"], stability["spearman_ST_rank"], color=[palette.get(x, "#999999") for x in stability["crop"]])
    ax.axhline(0.8, color="#C44E52", ls="--", lw=1)
    ax.set_ylim(0, 1.05)
    ax.set_ylabel("ST Spearman")
    ax.tick_params(axis="x", labelsize=7)
    fig.tight_layout()
    stab_path = ASSETS / "stability_spearman_en.png"
    fig.savefig(stab_path, dpi=320, bbox_inches="tight")
    plt.close(fig)
    return heat_path, gy_path, stab_path


def build():
    heat_path, gy_path, stab_path = make_figures()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    info = []

    def rec(t, b, c, n):
        info.append({"title": t, "bullets": b, "chart": c, "notes": n})

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, u(r"\u57fa\u4e8e APSIM Classic \u7684\u5c0f\u9ea6\u548c\u7389\u7c73\u54c1\u79cd\u53c2\u6570 Sobol \u654f\u611f\u6027\u5206\u6790"), u(r"APSIM Classic 7.10\uff5cWheat.xml / Maize.xml\uff5cN128 \u4e3b\u7ed3\u679c + N64 \u7a33\u5b9a\u6027\u9a8c\u8bc1"))
    for i, (v, lab, col) in enumerate([
        ("13", u(r"\u7b5b\u9009\u53c2\u6570"), C["blue"]),
        ("1920", u(r"\u6a21\u62df\u6b21\u6570"), C["green"]),
        ("5", u(r"\u76ee\u6807\u53d8\u91cf"), C["teal"]),
        ("0", u(r"\u8fd0\u884c\u9519\u8bef"), C["orange"]),
    ]):
        metric(slide, 0.75 + i * 3.0, 1.55, 2.35, 0.9, v, lab, col)
    b = [
        u(r"\u76ee\u6807\uff1a\u8bc6\u522b\u5c0f\u9ea6\u4e0e\u7389\u7c73\u54c1\u79cd\u53c2\u6570\u5bf9\u4ea7\u91cf\u3001\u7269\u5019\u3001LAI \u548c\u751f\u7269\u91cf\u7684\u4e3b\u5bfc\u5f71\u54cd\u3002"),
        u(r"\u65b9\u6cd5\uff1aSobol \u5168\u5c40\u654f\u611f\u6027\u5206\u6790\uff0c\u4ee5\u603b\u6548\u5e94\u6307\u6570 ST \u4e3a\u4e3b\u8981\u89e3\u91ca\u6307\u6807\u3002"),
        u(r"\u524d\u63d0\uff1aAPSIM Classic \u54c1\u79cd\u53c2\u6570\u6765\u81ea Wheat.xml \u548c Maize.xml\uff0c\u4e0d\u662f .apsimx\u3002"),
    ]
    bullets(slide, 0.85, 3.0, 11.5, 2.35, b, size=21)
    takeaway(slide, u(r"\u4e00\u53e5\u8bdd\uff1a\u7389\u7c73\u654f\u611f\u6027\u66f4\u96c6\u4e2d\u4e8e\u70ed\u65f6\u95f4/\u7269\u5019\uff0c\u5c0f\u9ea6\u5219\u7531\u7c7d\u7c92\u5927\u5c0f\u3001\u704c\u6d46\u4e0e\u7269\u5019\u5171\u540c\u63a7\u5236\u3002"))
    rec(u(r"\u6807\u9898\u4e0e\u7814\u7a76\u76ee\u6807"), b, u(r"\u6307\u6807\u5361\u7247"), u(r"\u5f00\u573a\u8bf4\u660e\u7814\u7a76\u5bf9\u8c61\u3001\u6a21\u578b\u7248\u672c\u548c Sobol \u5206\u6790\u76ee\u6807\u3002"))

    # 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, u(r"\u6a21\u578b\u8bbe\u7f6e\u4e0e\u6280\u672f\u8def\u7ebf"), u(r"\u4ece XML \u53c2\u6570\u63d0\u53d6\u5230 N64/N128 \u7a33\u5b9a\u6027\u9a8c\u8bc1"))
    steps = [u(r"\u54c1\u79cd\u8bc6\u522b"), u(r"\u53c2\u6570\u63d0\u53d6"), u(r"\u53c2\u6570\u7b5b\u9009"), u(r"Sobol \u62bd\u6837"), u(r"APSIM \u6279\u91cf\u8fd0\u884c"), u(r"\u8f93\u51fa\u63d0\u53d6"), u(r"\u6307\u6570\u8ba1\u7b97"), u(r"\u7a33\u5b9a\u6027\u9a8c\u8bc1")]
    for i, step in enumerate(steps):
        x = 0.55 + (i % 4) * 3.15
        y = 1.45 + (i // 4) * 1.65
        panel(slide, x, y, 2.55, 0.78)
        textbox(slide, x + 0.12, y + 0.22, 2.3, 0.28, step, size=17, bold=True, align=PP_ALIGN.CENTER)
        if i % 4 != 3:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.58), Inches(y + 0.25), Inches(0.38), Inches(0.25))
            arr.fill.solid()
            arr.fill.fore_color.rgb = C["line"]
            arr.line.color.rgb = C["line"]
    b = [
        u(r"\u5c0f\u9ea6\u54c1\u79cd\uff1aJimai70_v132_joint_iter353\uff1b\u7389\u7c73\u54c1\u79cd\uff1aP01_shandong_2025_v503_joint_iter60\u3002"),
        u(r"\u6bcf\u4e2a\u6837\u672c\u4fdd\u7559\u53c2\u6570\u8ffd\u8e2a\u8bb0\u5f55\uff1b\u4e34\u65f6\u66ff\u6362 crop XML \u540e\u6062\u590d baseline\uff0c\u5e76\u9010\u5b57\u8282\u6bd4\u5bf9\u3002"),
    ]
    bullets(slide, 0.85, 4.8, 11.6, 1.25, b, size=18)
    takeaway(slide, u(r"\u6d41\u7a0b\u8bbe\u8ba1\u76ee\u6807\uff1a\u53ef\u8ffd\u6eaf\u3001\u53ef\u6062\u590d\u3001\u53ef\u590d\u73b0\u3002"))
    rec(u(r"\u6a21\u578b\u8bbe\u7f6e\u4e0e\u6280\u672f\u8def\u7ebf"), b, u(r"\u539f\u751f\u6d41\u7a0b\u56fe"), u(r"\u8bf4\u660e Classic XML \u53c2\u6570\u6765\u6e90\u4e0e\u81ea\u52a8\u5316\u8fd0\u884c\u6d41\u7a0b\u3002"))

    # 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, u(r"\u8d28\u91cf\u63a7\u5236\u4e0e\u603b\u4f53\u654f\u611f\u6027\u683c\u5c40"), u(r"N128 \u8fd0\u884c\u5b8c\u6574\uff0cST heatmap \u663e\u793a\u4f5c\u7269\u548c\u8f93\u51fa\u53d8\u91cf\u5dee\u5f02"))
    metric(slide, 0.65, 1.15, 2.1, 0.85, "1920/1920", u(r"\u8fd0\u884c\u5b8c\u6210"), C["green"])
    metric(slide, 2.95, 1.15, 2.1, 0.85, "100%", u(r"\u6210\u529f\u7387"), C["blue"])
    metric(slide, 5.25, 1.15, 2.1, 0.85, "no diff", u(r"XML \u6062\u590d"), C["teal"])
    picture(slide, heat_path, 0.75, 2.25, w=7.05)
    b = [
        u(r"\u7389\u7c73\u4ea7\u91cf\u548c\u7269\u5019\u4e3b\u8981\u96c6\u4e2d\u5728 thermal-time / grain-filling \u53c2\u6570\u3002"),
        u(r"\u5c0f\u9ea6\u4ea7\u91cf\u654f\u611f\u6027\u66f4\u5206\u6563\uff0c\u6d89\u53ca grain size\u3001grain filling\u3001photoperiod \u548c vernalization\u3002"),
        u(r"\u4e0d\u540c\u76ee\u6807\u53d8\u91cf\u7684\u4e3b\u5bfc\u53c2\u6570\u4e0d\u540c\uff0c\u6821\u51c6\u76ee\u6807\u4f1a\u5f71\u54cd\u53c2\u6570\u4f18\u5148\u7ea7\u3002"),
    ]
    bullets(slide, 8.2, 2.25, 4.35, 2.9, b, size=18)
    takeaway(slide, u(r"\u603b\u4f53\u7ed3\u679c\u652f\u6301\u201c\u4f5c\u7269\u7279\u5f02\u3001\u8f93\u51fa\u7279\u5f02\u201d\u7684\u54c1\u79cd\u53c2\u6570\u654f\u611f\u6027\u7ed3\u6784\u3002"))
    rec(u(r"\u8d28\u91cf\u63a7\u5236\u4e0e\u603b\u4f53\u654f\u611f\u6027\u683c\u5c40"), b, u(r"ST heatmap + QC \u6307\u6807"), u(r"\u5148\u8bf4\u660e\u8fd0\u884c\u53ef\u9760\uff0c\u518d\u8fdb\u5165\u603b\u4f53\u654f\u611f\u6027\u683c\u5c40\u3002"))

    # 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, u(r"\u4ea7\u91cf\u654f\u611f\u6027\uff1a\u7389\u7c73\u96c6\u4e2d\uff0c\u5c0f\u9ea6\u5206\u6563"), u(r"grain_yield \u7684 Top ST \u53c2\u6570\u63ed\u793a\u4e0d\u540c\u4f5c\u7269\u673a\u5236"))
    picture(slide, gy_path, 0.7, 1.2, w=7.0)
    b = [
        u(r"\u7389\u7c73 Top \u53c2\u6570\uff1att_flag_to_flower\u3001tt_flower_to_maturity\u3001tt_endjuv_to_init\u3002"),
        u(r"\u89e3\u91ca\uff1a\u5f00\u82b1\u65f6\u5e8f\u548c\u704c\u6d46/\u6210\u719f\u957f\u5ea6\u6539\u53d8\u751f\u7269\u91cf\u79ef\u7d2f\u7a97\u53e3\u3002"),
        u(r"\u5c0f\u9ea6 Top \u53c2\u6570\uff1amax_grain_size\u3001tt_end_of_juvenile\u3001tt_floral_initiation\u3001tt_start_grain_fill\u3001vern_sens\u3002"),
        u(r"\u89e3\u91ca\uff1a\u5c0f\u9ea6\u4ea7\u91cf\u53d7 sink potential\u3001grain filling \u4e0e phenology \u5171\u540c\u5f71\u54cd\u3002"),
    ]
    bullets(slide, 8.05, 1.2, 4.55, 4.2, b, size=17)
    takeaway(slide, u(r"\u7389\u7c73\u66f4\u50cf\u5355\u4e00\u8def\u5f84\u4e3b\u5bfc\uff1b\u5c0f\u9ea6\u9700\u8981\u7528\u591a\u53c2\u6570\u7ec4\u5171\u540c\u4f5c\u7528\u6765\u89e3\u91ca\u3002"))
    rec(u(r"\u4ea7\u91cf\u654f\u611f\u6027"), b, u(r"\u7389\u7c73\u548c\u5c0f\u9ea6 grain_yield Top ST \u6761\u5f62\u56fe"), u(r"\u8fd9\u662f\u7ed3\u679c\u6838\u5fc3\u9875\uff0c\u91cd\u70b9\u89e3\u91ca\u53c2\u6570\u80cc\u540e\u7684\u751f\u7406\u542b\u4e49\u3002"))

    # 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, u(r"\u7a33\u5b9a\u6027\u9a8c\u8bc1\u4e0e\u4e0b\u4e00\u6b65\u5de5\u4f5c"), u(r"N64 \u4e0e N128 \u7684 ST \u6392\u540d\u9ad8\u5ea6\u4e00\u81f4"))
    picture(slide, stab_path, 0.75, 1.15, w=6.4)
    for i, (val, lab, col) in enumerate([
        ("10/10", "Top3 overlap \u2265 2", C["green"]),
        ("10/10", "Top5 overlap \u2265 4", C["green"]),
        ("10/10", "ST Spearman \u2265 0.8", C["green"]),
        ("0/10", "ST Spearman < 0.6", C["orange"]),
    ]):
        metric(slide, 7.75, 1.1 + i * 1.03, 4.0, 0.72, val, lab, col)
    b = [
        u(r"\u7ed3\u8bba\uff1aN128 \u53ef\u4f5c\u4e3a\u4e3b\u7ed3\u679c\uff0cN64 \u4f5c\u4e3a\u7a33\u5b9a\u6027\u9a8c\u8bc1\uff1bS1 \u4ec5\u4f5c\u8f85\u52a9\u8bf4\u660e\u3002"),
        u(r"\u9650\u5236\u4e0e\u4e0b\u4e00\u6b65\uff1agrain_number\u3001grain_weight\u3001water_use_efficiency \u5f53\u524d\u672a\u8f93\u51fa\uff0c\u9700\u5148\u5b8c\u5584 APSIM output/report \u5e76\u7528 --limit 5 \u5c0f\u6837\u672c\u6d4b\u8bd5\u3002"),
    ]
    bullets(slide, 0.9, 5.15, 11.6, 1.0, b, size=17)
    takeaway(slide, u(r"\u6700\u7ec8\u603b\u7ed3\uff1a\u5f53\u524d\u7ed3\u679c\u5df2\u7ecf\u652f\u6301\u4f5c\u7269\u5dee\u5f02\u5316\u7684 APSIM \u54c1\u79cd\u53c2\u6570\u654f\u611f\u6027\u4e3b\u7ebf\u3002"))
    rec(u(r"\u7a33\u5b9a\u6027\u9a8c\u8bc1\u4e0e\u4e0b\u4e00\u6b65"), b, u(r"ST Spearman \u7a33\u5b9a\u6027\u67f1\u72b6\u56fe + \u6307\u6807\u5361\u7247"), u(r"\u7528\u7a33\u5b9a\u6027\u7ed3\u679c\u6536\u675f\uff0c\u8bf4\u660e\u4e0b\u4e00\u6b65\u6269\u5c55\u4e0d\u5f71\u54cd\u5f53\u524d\u4e3b\u7ed3\u8bba\u3002"))

    ppt = OUT / "APSIM_Classic_Sobol_CN_5slides_UTF8_verified.pptx"
    prs.save(ppt)
    notes_lines = [u(r"# 5\u9875\u538b\u7f29\u7248\u4e2d\u6587\u6c47\u62a5 PPT \u5185\u5bb9\u4e0e\u8bb2\u7a3f\u5907\u6ce8"), ""]
    for i, item in enumerate(info, 1):
        notes_lines += [f"## 第 {i} 页：{item['title']}", "", "**Bullet points**"]
        notes_lines += [f"- {x}" for x in item["bullets"]]
        notes_lines += ["", f"**建议图表**：{item['chart']}", "", f"**讲稿备注**：{item['notes']}", ""]
    notes = OUT / "slide_titles_bullets_notes_cn_5slides.md"
    notes.write_text("\n".join(notes_lines), encoding="utf-8")
    return ppt, notes


def verify(ppt: Path):
    with zipfile.ZipFile(ppt, "r") as z:
        slides = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        xml = "\n".join(z.read(n).decode("utf-8", errors="strict") for n in slides)
        media = len([n for n in z.namelist() if n.startswith("ppt/media/")])
    result = {
        "slide_count": len(slides),
        "media_count": media,
        "chinese_char_count": len(re.findall(r"[\u4e00-\u9fff]", xml)),
        "has_question_damage": "????" in xml,
        "font_count": xml.count(FONT),
    }
    if result["has_question_damage"] or result["chinese_char_count"] <= 0 or result["slide_count"] != 5:
        raise RuntimeError(json.dumps(result, ensure_ascii=False))
    return result


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    ppt, notes = build()
    result = verify(ppt)
    qa = OUT / "qa_report_5slides.md"
    qa.write_text(
        "\n".join([
            "# QA report",
            "",
            f"- PPTX: `{ppt}`",
            f"- Notes: `{notes}`",
            f"- Slide count: {result['slide_count']}",
            f"- Embedded media count: {result['media_count']}",
            f"- Chinese character count in slide XML: {result['chinese_char_count']}",
            f"- Contains question-mark damage: {result['has_question_damage']}",
            f"- Microsoft YaHei occurrences: {result['font_count']}",
            "",
            "All Chinese text is native PowerPoint text with East Asian font set to Microsoft YaHei.",
        ]),
        encoding="utf-8",
    )
    print(json.dumps(result, ensure_ascii=False))
    print(f"PPTX={ppt}")
    print(f"NOTES={notes}")
    print(f"QA={qa}")


if __name__ == "__main__":
    main()
