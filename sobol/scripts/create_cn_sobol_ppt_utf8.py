# -*- coding: utf-8 -*-
"""
Create a Chinese PPTX report for APSIM Classic Sobol sensitivity analysis.

Encoding/font safety rules:
- Python source is UTF-8.
- All Chinese text is inserted as native PowerPoint text runs.
- Each run explicitly sets both run.font.name and East Asian font to Microsoft YaHei.
- The script verifies PPTX XML contains real Chinese characters and does not contain "????".
"""

from __future__ import annotations

import json
import re
import zipfile
from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


BASE = Path(r"F:\APSIM710-r4221\process_bio\sobol")
N128_FINAL = BASE / "organized_outputs_screened_N128_20260515_185604" / "final_results"
OUT_DIR = BASE / "ppt_report_sobol_20260517_utf8"
ASSET_DIR = OUT_DIR / "assets" / "figures"
FONT_FACE = "Microsoft YaHei"


def u(text: str) -> str:
    """Decode escaped Unicode text into real Python Unicode."""
    return text.encode("ascii").decode("unicode_escape")


ZH_TEST = u(r"\u4e2d\u6587\u6d4b\u8bd5\uff1a\u5c0f\u9ea6\u3001\u7389\u7c73\u3001\u54c1\u79cd\u53c2\u6570\u3001\u654f\u611f\u6027\u5206\u6790")


def read_csv_utf8(path: Path) -> pd.DataFrame:
    try:
        return pd.read_csv(path, encoding="utf-8-sig")
    except UnicodeDecodeError:
        return pd.read_csv(path, encoding="gb18030")


def set_run_font(run, size: int = 18, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.name = FONT_FACE
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        elem = rpr.find(qn(tag))
        if elem is None:
            elem = OxmlElement(tag)
            rpr.append(elem)
        elem.set("typeface", FONT_FACE)


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 18,
    bold: bool = False,
    color: RGBColor | None = None,
    align=None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    set_run_font(r, size=size, bold=bold, color=color or COLORS["navy"])
    return box


def add_bullets(slide, x: float, y: float, w: float, h: float, bullets: list[str], size: int = 19):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.06)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.space_after = Pt(7)
        for run in p.runs:
            set_run_font(run, size=size, color=COLORS["navy"])
    return box


def add_title(slide, title: str, subtitle: str | None = None):
    add_textbox(slide, 0.45, 0.25, 12.4, 0.62, title, size=30, bold=True, color=COLORS["navy"])
    if subtitle:
        add_textbox(slide, 0.48, 0.88, 12.2, 0.34, subtitle, size=13, color=COLORS["gray"])


def add_takeaway(slide, text: str):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.77), Inches(12.25), Inches(0.44))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(235, 244, 242)
    shp.line.color.rgb = RGBColor(194, 218, 214)
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run_font(r, size=15, bold=True, color=COLORS["navy"])


def add_metric(slide, x, y, w, h, value: str, label: str, color: RGBColor):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(248, 250, 251)
    shp.line.color.rgb = color
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    set_run_font(r, size=25, bold=True, color=color)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    set_run_font(r2, size=12, color=COLORS["gray"])


def add_panel(slide, x, y, w, h, fill=RGBColor(248, 250, 251)):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = COLORS["line"]
    return shp


def add_picture(slide, path: Path, x, y, w=None, h=None):
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if w is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    if h is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def add_caption(slide, x, y, w, text: str):
    add_textbox(slide, x, y, w, 0.28, text, size=13, color=COLORS["gray"])


def add_source(slide, text: str):
    add_textbox(slide, 0.55, 7.15, 12.25, 0.2, text, size=9, color=COLORS["gray"])


COLORS = {
    "navy": RGBColor(28, 51, 74),
    "blue": RGBColor(61, 104, 153),
    "teal": RGBColor(70, 145, 141),
    "green": RGBColor(78, 132, 83),
    "orange": RGBColor(196, 116, 65),
    "red": RGBColor(180, 70, 72),
    "gray": RGBColor(88, 96, 104),
    "line": RGBColor(210, 218, 226),
}


def make_figures() -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    idx = read_csv_utf8(N128_FINAL / "sobol_indices_summary.csv")
    for col in ["S1", "S1_conf", "ST", "ST_conf"]:
        idx[col] = pd.to_numeric(idx[col], errors="coerce")
    idx = idx[idx["ST"].notna()].copy()
    idx["parameter_name"] = idx["parameter_key"].astype(str).str.split("__").str[2]
    idx["parameter_name"] = idx["parameter_name"].str.replace("_1", "[1]", regex=False).str.replace("_2", "[2]", regex=False)

    outputs = read_csv_utf8(N128_FINAL / "sobol_model_outputs.csv")
    stability = read_csv_utf8(N128_FINAL / "sobol_N64_vs_N128_stability_report.csv")
    stability = stability[["crop", "target_variable", "spearman_ST_rank"]].drop_duplicates()

    sns.set_theme(style="whitegrid", context="talk")
    plt.rcParams["font.family"] = "DejaVu Sans"
    palette = {"maize": "#4C78A8", "wheat": "#72B7B2"}

    fig, axes = plt.subplots(1, 3, figsize=(10.5, 3.0))
    for ax, var in zip(axes, ["grain_yield", "biomass", "lai"]):
        df = outputs[["crop", var]].rename(columns={var: "value"})
        sns.boxplot(data=df, x="crop", y="value", order=["maize", "wheat"], palette=[palette["maize"], palette["wheat"]], ax=ax)
        ax.set_title(var, fontsize=12)
        ax.set_xlabel("")
        ax.set_ylabel("")
    fig.tight_layout()
    output_dist = ASSET_DIR / "output_distributions_en.png"
    fig.savefig(output_dist, dpi=320, bbox_inches="tight")
    plt.close(fig)

    heat = idx.pivot_table(index="parameter_name", columns=["crop", "target_variable"], values="ST", aggfunc="mean").fillna(0)
    heat = heat.loc[heat.max(axis=1).sort_values(ascending=False).index]
    fig, ax = plt.subplots(figsize=(8.5, 5.2))
    sns.heatmap(heat, cmap="YlGnBu", linewidths=0.25, linecolor="white", cbar_kws={"label": "ST"}, ax=ax)
    ax.set_xlabel("crop / target")
    ax.set_ylabel("parameter")
    ax.tick_params(axis="x", labelrotation=45, labelsize=8)
    ax.tick_params(axis="y", labelsize=8)
    fig.tight_layout()
    heatmap = ASSET_DIR / "st_heatmap_en.png"
    fig.savefig(heatmap, dpi=320, bbox_inches="tight")
    plt.close(fig)

    bar_paths = {}
    for crop in ["maize", "wheat"]:
        df = idx[(idx["crop"] == crop) & (idx["target_variable"] == "grain_yield")].sort_values("ST", ascending=False).head(6)
        fig, ax = plt.subplots(figsize=(5.4, 3.6))
        ax.barh(df["parameter_name"][::-1], df["ST"][::-1], color=palette[crop])
        ax.set_xlabel("Sobol total-effect index (ST)")
        ax.set_ylabel("")
        ax.set_xlim(0, max(0.95 if crop == "maize" else 0.38, df["ST"].max() * 1.18))
        for i, val in enumerate(df["ST"][::-1]):
            ax.text(val + 0.01, i, f"{val:.2f}", va="center", fontsize=9)
        fig.tight_layout()
        path = ASSET_DIR / f"{crop}_grain_yield_top_st_en.png"
        fig.savefig(path, dpi=320, bbox_inches="tight")
        plt.close(fig)
        bar_paths[crop] = path

    stability["group"] = stability["crop"] + "\n" + stability["target_variable"]
    fig, ax = plt.subplots(figsize=(8.0, 3.4))
    ax.bar(stability["group"], stability["spearman_ST_rank"], color=[palette.get(x, "#999999") for x in stability["crop"]])
    ax.axhline(0.8, color="#C44E52", ls="--", lw=1)
    ax.set_ylim(0, 1.05)
    ax.set_ylabel("ST Spearman")
    ax.tick_params(axis="x", labelsize=7)
    fig.tight_layout()
    stable = ASSET_DIR / "stability_spearman_en.png"
    fig.savefig(stable, dpi=320, bbox_inches="tight")
    plt.close(fig)

    return {
        "output_dist": output_dist,
        "heatmap": heatmap,
        "maize_bar": bar_paths["maize"],
        "wheat_bar": bar_paths["wheat"],
        "stability": stable,
    }


def build_pptx(figs: dict[str, Path]) -> tuple[Path, list[dict[str, object]]]:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_info: list[dict[str, object]] = []

    def record(title, bullets, chart, notes):
        slide_info.append({"title": title, "bullets": bullets, "chart": chart, "notes": notes})

    # Test slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, ZH_TEST, u(r"\u5982\u672c\u9875\u4e0d\u80fd\u6b63\u5e38\u663e\u793a\u4e2d\u6587\uff0c\u5219\u4e0d\u5e94\u7ee7\u7eed\u4f7f\u7528\u8be5 PPT\u3002"))
    add_bullets(slide, 0.9, 1.6, 11.5, 2.0, [ZH_TEST], size=24)
    add_takeaway(slide, u(r"\u4e2d\u6587\u663e\u793a\u6d4b\u8bd5\u9875\uff1a\u8be5\u9875\u7528\u4e8e\u68c0\u67e5 PPTX \u751f\u6210\u9636\u6bb5\u662f\u5426\u4fdd\u7559\u771f\u5b9e\u6c49\u5b57\u3002"))

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, u(r"\u57fa\u4e8e APSIM Classic \u7684\u5c0f\u9ea6\u548c\u7389\u7c73\u54c1\u79cd\u53c2\u6570 Sobol \u654f\u611f\u6027\u5206\u6790"), u(r"\u6a21\u578b\uff1aAPSIM Classic 7.10\uff5c\u4f5c\u7269\uff1a\u5c0f\u9ea6\u4e0e\u7389\u7c73\uff5c\u65b9\u6cd5\uff1aSALib Sobol \u5168\u5c40\u654f\u611f\u6027\u5206\u6790"))
    add_textbox(slide, 0.75, 1.45, 4.5, 0.4, u(r"\u6c47\u62a5\u4eba\uff1aXXX"), size=20)
    add_textbox(slide, 0.75, 1.9, 4.5, 0.4, u(r"\u65e5\u671f\uff1a2026\u5e745\u670817\u65e5"), size=20)
    add_textbox(slide, 0.75, 2.35, 8.7, 0.4, u(r"\u7814\u7a76\u5bf9\u8c61\uff1aJimai70_v132_joint_iter353\uff1bP01_shandong_2025_v503_joint_iter60"), size=18, color=COLORS["gray"])
    for i, (value, label, color) in enumerate([
        ("13", u(r"\u7b5b\u9009\u540e\u53c2\u6570"), COLORS["blue"]),
        ("128", "Sobol N", COLORS["teal"]),
        ("1920", u(r"\u6a21\u62df\u6b21\u6570"), COLORS["green"]),
        ("100%", u(r"\u8fd0\u884c\u6210\u529f\u7387"), COLORS["orange"]),
    ]):
        add_metric(slide, 0.85 + i * 3.0, 3.35, 2.45, 1.0, value, label, color)
    add_takeaway(slide, u(r"\u6838\u5fc3\u7ed3\u8bba\uff1aN128 \u7ed3\u679c\u7a33\u5b9a\uff0c\u7389\u7c73\u4ee5\u70ed\u65f6\u95f4/\u7269\u5019\u53c2\u6570\u4e3b\u5bfc\uff0c\u5c0f\u9ea6\u8868\u73b0\u4e3a\u7c7d\u7c92\u5927\u5c0f\u3001\u704c\u6d46\u4e0e\u7269\u5019\u5171\u540c\u63a7\u5236\u3002"))
    add_source(slide, r"F:\APSIM710-r4221\process_bio\sobol")
    record(u(r"\u6807\u9898\u9875"), [u(r"\u4e3b\u9898\u3001\u6c47\u62a5\u4eba\u3001\u65e5\u671f\u3001\u7814\u7a76\u5bf9\u8c61"), u(r"\u7528 13 \u4e2a\u7b5b\u9009\u53c2\u6570\u5b8c\u6210 N128 \u4e3b\u5206\u6790")], u(r"\u6307\u6807\u5361\u7247"), u(r"\u5f00\u573a\u8bf4\u660e\u8fd9\u662f\u57fa\u4e8e\u5df2\u5b8c\u6210 N128 \u4e3b\u5206\u6790\u7684\u9636\u6bb5\u6027\u6c47\u62a5\u3002"))

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, u(r"\u7814\u7a76\u80cc\u666f\u4e0e\u76ee\u6807"), u(r"\u4e3a\u4ec0\u4e48\u8981\u505a\u54c1\u79cd\u53c2\u6570\u7684\u5168\u5c40\u654f\u611f\u6027\u5206\u6790\uff1f"))
    bullets = [
        u(r"APSIM \u54c1\u79cd\u53c2\u6570\u51b3\u5b9a\u7269\u5019\u3001\u704c\u6d46\u3001\u51a0\u5c42\u548c\u4ea7\u91cf\u5f62\u6210\u8fc7\u7a0b\u3002"),
        u(r"\u53c2\u6570\u591a\u4e14\u53ef\u80fd\u5b58\u5728\u975e\u7ebf\u6027\u548c\u4ea4\u4e92\uff0c\u5355\u56e0\u7d20\u6270\u52a8\u96be\u4ee5\u8bc6\u522b\u5168\u5c40\u8d21\u732e\u3002"),
        u(r"Sobol \u65b9\u6cd5\u53ef\u5206\u89e3\u8f93\u51fa\u65b9\u5dee\uff0c\u5f97\u5230\u4e00\u9636\u6548\u5e94 S1 \u548c\u603b\u6548\u5e94 ST\u3002"),
        u(r"\u76ee\u6807\uff1a\u8bc6\u522b\u5c0f\u9ea6\u4e0e\u7389\u7c73\u54c1\u79cd\u53c2\u6570\u5bf9\u5173\u952e\u8f93\u51fa\u7684\u4e3b\u5bfc\u5f71\u54cd\u8def\u5f84\u3002"),
    ]
    add_bullets(slide, 0.75, 1.35, 5.8, 4.7, bullets, size=21)
    add_panel(slide, 7.1, 1.35, 4.9, 3.1)
    add_textbox(slide, 7.45, 1.65, 4.2, 0.45, u(r"\u5206\u6790\u76ee\u6807"), size=24, bold=True, color=COLORS["blue"])
    add_bullets(slide, 7.45, 2.25, 4.1, 1.8, [u(r"\u7b5b\u51fa\u5173\u952e\u54c1\u79cd\u53c2\u6570"), u(r"\u89e3\u91ca\u4f5c\u7269\u95f4\u654f\u611f\u6027\u5dee\u5f02"), u(r"\u4e3a APSIM cultivar calibration \u63d0\u4f9b\u4f18\u5148\u7ea7")], size=18)
    add_takeaway(slide, u(r"Sobol ST \u662f\u672c\u6c47\u62a5\u7684\u4e3b\u8981\u89e3\u91ca\u6307\u6807\uff1bS1 \u4ec5\u4f5c\u4e3a\u76f4\u63a5\u4e3b\u6548\u5e94\u7684\u8f85\u52a9\u53c2\u8003\u3002"))
    record(u(r"\u7814\u7a76\u80cc\u666f\u4e0e\u76ee\u6807"), bullets, u(r"\u6982\u5ff5\u6846\uff1a\u5206\u6790\u76ee\u6807"), u(r"\u5f3a\u8c03\u4f7f\u7528\u5168\u5c40 Sobol \u662f\u56e0\u4e3a\u53c2\u6570\u53ef\u80fd\u5b58\u5728\u975e\u7ebf\u6027\u548c\u4ea4\u4e92\u3002"))

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, u(r"\u6a21\u578b\u4e0e\u6570\u636e\u8bbe\u7f6e"), u(r"APSIM Classic cultivar \u53c2\u6570\u6765\u81ea Wheat.xml \u548c Maize.xml"))
    cards = [
        (u(r"\u6a21\u578b\u7248\u672c"), u(r"APSIM Classic 7.10\n\u4e0d\u662f APSIM Next Generation"), COLORS["blue"]),
        (u(r"\u53c2\u6570\u6765\u6e90"), u(r"Wheat.xml / Maize.xml\n\u4e0d\u662f .apsimx JSON"), COLORS["teal"]),
        (u(r"\u5206\u6790\u54c1\u79cd"), u(r"\u5c0f\u9ea6\uff1aJimai70_v132_joint_iter353\n\u7389\u7c73\uff1aP01_shandong_2025_v503_joint_iter60"), COLORS["green"]),
    ]
    for i, (head, body, color) in enumerate(cards):
        add_panel(slide, 0.75 + i * 4.1, 1.35, 3.55, 1.65)
        add_textbox(slide, 1.0 + i * 4.1, 1.55, 3.1, 0.35, head, size=22, bold=True, color=color)
        add_textbox(slide, 1.0 + i * 4.1, 2.02, 3.15, 0.75, body, size=16)
    add_textbox(slide, 0.85, 3.65, 5.5, 0.35, u(r"\u53ef\u89e3\u91ca\u8f93\u51fa\u53d8\u91cf"), size=22, bold=True)
    add_bullets(slide, 1.0, 4.15, 4.7, 1.5, ["grain_yield", "biomass", "lai", "flowering_date", "maturity_date"], size=18)
    add_textbox(slide, 7.0, 3.65, 5.3, 0.35, u(r"\u6682\u672a\u7eb3\u5165\u53d8\u91cf"), size=22, bold=True, color=COLORS["red"])
    add_bullets(slide, 7.15, 4.15, 4.9, 1.2, ["grain_number", "grain_weight", "water_use_efficiency"], size=18)
    add_takeaway(slide, u(r"\u5173\u952e\u6280\u672f\u70b9\uff1a\u5fc5\u987b\u6309 APSIM Classic crop XML \u7ed3\u6784\u8bfb\u53d6\u548c\u4fee\u6539 cultivar \u53c2\u6570\u3002"))
    record(u(r"\u6a21\u578b\u4e0e\u6570\u636e\u8bbe\u7f6e"), [u(r"APSIM Classic 7.10"), u(r"\u53c2\u6570\u6765\u81ea Wheat.xml \u548c Maize.xml"), u(r"5 \u4e2a\u8f93\u51fa\u53d8\u91cf\u53ef\u89e3\u91ca\uff0c3 \u4e2a\u53d8\u91cf\u6682\u672a\u8f93\u51fa")], u(r"\u4fe1\u606f\u5361\u7247\u4e0e\u53d8\u91cf\u6e05\u5355"), u(r"\u63d0\u9192\u542c\u4f17 Classic \u548c Next Gen \u6587\u4ef6\u7ed3\u6784\u4e0d\u540c\u3002"))

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, u(r"\u6280\u672f\u8def\u7ebf"), u(r"\u4ece\u54c1\u79cd\u8bc6\u522b\u5230 N64/N128 \u7a33\u5b9a\u6027\u9a8c\u8bc1"))
    steps = [u(r"\u54c1\u79cd\u8bc6\u522b"), u(r"\u53c2\u6570\u63d0\u53d6"), u(r"\u53c2\u6570\u7b5b\u9009"), u(r"Sobol \u62bd\u6837"), u(r"APSIM \u6279\u91cf\u8fd0\u884c"), u(r"\u8f93\u51fa\u63d0\u53d6"), u(r"Sobol \u6307\u6570\u8ba1\u7b97"), u(r"\u7a33\u5b9a\u6027\u9a8c\u8bc1")]
    for i, step in enumerate(steps):
        x = 0.55 + (i % 4) * 3.15
        y = 1.55 + (i // 4) * 2.05
        add_panel(slide, x, y, 2.55, 0.95, RGBColor(245, 249, 250))
        add_textbox(slide, x + 0.15, y + 0.27, 2.25, 0.3, step, size=18, bold=True, align=PP_ALIGN.CENTER)
        if i % 4 != 3:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.58), Inches(y + 0.32), Inches(0.42), Inches(0.28))
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLORS["line"]
            arr.line.color.rgb = COLORS["line"]
    add_bullets(slide, 0.85, 5.75, 11.7, 0.65, [u(r"\u6bcf\u4e2a\u6837\u672c\u5747\u4fdd\u7559\u53c2\u6570\u8ffd\u8e2a\u8bb0\u5f55\uff1b\u6bcf\u6b21\u4e34\u65f6\u66ff\u6362 crop XML \u540e\u6062\u590d baseline\uff0c\u5e76\u8fdb\u884c\u6587\u4ef6\u6bd4\u5bf9\u3002")], size=18)
    add_takeaway(slide, u(r"\u6d41\u7a0b\u8bbe\u8ba1\u91cd\u70b9\uff1a\u53ef\u8ffd\u6eaf\u3001\u53ef\u6062\u590d\u3001\u53ef\u590d\u73b0\u3002"))
    record(u(r"\u6280\u672f\u8def\u7ebf"), [u(r"\u54c1\u79cd\u8bc6\u522b \u2192 \u53c2\u6570\u63d0\u53d6 \u2192 \u53c2\u6570\u7b5b\u9009 \u2192 Sobol \u62bd\u6837"), u(r"APSIM \u6279\u91cf\u8fd0\u884c \u2192 \u8f93\u51fa\u63d0\u53d6 \u2192 \u6307\u6570\u8ba1\u7b97 \u2192 \u7a33\u5b9a\u6027\u9a8c\u8bc1")], u(r"\u539f\u751f PowerPoint \u6d41\u7a0b\u56fe"), u(r"\u8bf4\u660e\u81ea\u52a8\u5316\u6d41\u7a0b\u5982\u4f55\u4fdd\u8bc1\u6bcf\u4e2a\u6837\u672c\u53ef\u8ffd\u6eaf\u3002"))

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, u(r"Sobol \u8bbe\u8ba1\u4e0e\u8d28\u91cf\u63a7\u5236"), u(r"N128 \u4e3b\u5206\u6790\u5b8c\u6210 1920 \u6b21 APSIM Classic \u6a21\u62df"))
    for i, (value, label, color) in enumerate([
        ("D = 13", u(r"\u7b5b\u9009\u53c2\u6570\u6570"), COLORS["blue"]),
        ("N = 128", u(r"\u57fa\u7840\u6837\u672c\u91cf"), COLORS["teal"]),
        ("1920", u(r"\u603b\u6a21\u62df\u6570"), COLORS["green"]),
        ("0", u(r"\u8fd0\u884c\u9519\u8bef"), COLORS["orange"]),
        ("5", u(r"\u53ef\u89e3\u91ca\u53d8\u91cf"), COLORS["navy"]),
    ]):
        add_metric(slide, 0.65 + i * 2.5, 1.1, 2.1, 0.9, value, label, color)
    add_picture(slide, figs["output_dist"], 0.85, 2.45, w=6.8)
    add_bullets(slide, 8.15, 2.5, 4.35, 2.4, [u(r"APSIM \u8fd0\u884c\u72b6\u6001\uff1a1920 / 1920 finished\u3002"), u(r"Wheat.xml / Maize.xml \u6062\u590d\u68c0\u67e5\uff1ano differences encountered\u3002"), u(r"\u8f93\u51fa\u53d8\u91cf\u5b8c\u6574\uff1a\u4ea7\u91cf\u3001\u751f\u7269\u91cf\u3001LAI\u3001\u5f00\u82b1\u671f\u3001\u6210\u719f\u671f\u3002")], size=18)
    add_caption(slide, 0.95, 5.85, 6.3, u(r"\u56fe\uff1aN128 \u6837\u672c\u4e0b grain_yield\u3001biomass \u548c LAI \u7684\u8f93\u51fa\u5206\u5e03\u3002"))
    add_takeaway(slide, u(r"N128 \u6570\u636e\u8d28\u91cf\u6ee1\u8db3\u540e\u7eed\u654f\u611f\u6027\u89e3\u91ca\u548c\u8bba\u6587\u56fe\u8868\u6574\u7406\u8981\u6c42\u3002"))
    record(u(r"Sobol \u8bbe\u8ba1\u4e0e\u8d28\u91cf\u63a7\u5236"), [u(r"D=13\uff0cN=128\uff0c\u603b\u6a21\u62df\u6570=1920"), u(r"\u8fd0\u884c\u6210\u529f\u7387 100%\uff0c\u9519\u8bef\u6570 0"), u(r"XML \u6062\u590d\u68c0\u67e5\u65e0\u5dee\u5f02"), u(r"5 \u4e2a\u76ee\u6807\u53d8\u91cf\u5b8c\u6574")], u(r"\u8f93\u51fa\u53d8\u91cf\u5206\u5e03\u56fe + \u6307\u6807\u5361\u7247"), u(r"\u8fd9\u4e00\u9875\u7528\u4e8e\u8bf4\u660e\u8fd0\u884c\u548c\u8f93\u51fa\u90fd\u662f\u53ef\u68c0\u67e5\u7684\u3002"))

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, u(r"\u603b\u4f53\u654f\u611f\u6027\u7ed3\u679c"), u(r"\u4e0d\u540c\u4f5c\u7269\u548c\u76ee\u6807\u53d8\u91cf\u5448\u73b0\u4e0d\u540c\u7684 ST \u654f\u611f\u6027\u7ed3\u6784"))
    add_picture(slide, figs["heatmap"], 0.75, 1.15, w=7.1)
    add_bullets(slide, 8.25, 1.35, 4.25, 3.8, [u(r"\u7389\u7c73\u4ea7\u91cf\u548c\u7269\u5019\u4e3b\u8981\u96c6\u4e2d\u5728 thermal-time / grain-filling \u53c2\u6570\u3002"), u(r"\u5c0f\u9ea6\u4ea7\u91cf\u654f\u611f\u6027\u66f4\u5206\u6563\uff0c\u6d89\u53ca grain size\u3001grain filling\u3001photoperiod \u548c vernalization\u3002"), u(r"\u4e0d\u540c\u8f93\u51fa\u53d8\u91cf\u7684\u4e3b\u5bfc\u53c2\u6570\u4e0d\u540c\uff0c\u8bf4\u660e\u6821\u51c6\u76ee\u6807\u4f1a\u5f71\u54cd\u53c2\u6570\u4f18\u5148\u7ea7\u3002")], size=18)
    add_caption(slide, 0.85, 6.05, 7.1, u(r"\u56fe\uff1aN128 Sobol ST heatmap\uff0c\u989c\u8272\u8d8a\u6df1\u8868\u793a\u603b\u6548\u5e94\u8d8a\u5f3a\u3002"))
    add_takeaway(slide, u(r"\u603b\u4f53\u7ed3\u679c\u652f\u6301\u201c\u4f5c\u7269\u7279\u5f02\u3001\u8f93\u51fa\u7279\u5f02\u201d\u7684\u54c1\u79cd\u53c2\u6570\u654f\u611f\u6027\u7ed3\u6784\u3002"))
    record(u(r"\u603b\u4f53\u654f\u611f\u6027\u7ed3\u679c"), [u(r"ST heatmap \u663e\u793a\u4e0d\u540c crop \u00d7 target \u7684\u4e3b\u5bfc\u53c2\u6570\u4e0d\u540c"), u(r"\u7389\u7c73\u96c6\u4e2d\u4e8e\u70ed\u65f6\u95f4\u548c\u704c\u6d46"), u(r"\u5c0f\u9ea6\u6d89\u53ca\u66f4\u591a\u53c2\u6570\u7ec4")], u(r"ST heatmap"), u(r"\u63d0\u9192\u542c\u4f17 ST \u662f\u603b\u6548\u5e94\uff0c\u9002\u5408\u4f5c\u4e3a\u4e3b\u89e3\u91ca\u6307\u6807\u3002"))

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, u(r"\u7389\u7c73\u4ea7\u91cf\u654f\u611f\u6027\u7ed3\u679c"), u(r"grain_yield \u4e3b\u8981\u53d7\u5f00\u82b1\u548c\u6210\u719f\u76f8\u5173\u70ed\u65f6\u95f4\u53c2\u6570\u63a7\u5236"))
    add_picture(slide, figs["maize_bar"], 0.8, 1.25, w=5.8)
    bullets = [u(r"tt_flag_to_flower \u662f\u6700\u5f3a\u654f\u611f\u53c2\u6570\uff0cST = 0.868\u3002"), u(r"tt_flower_to_maturity \u5f71\u54cd\u5f00\u82b1\u5230\u6210\u719f\u9636\u6bb5\uff0c\u5173\u7cfb\u5230\u704c\u6d46\u6301\u7eed\u65f6\u95f4\u3002"), u(r"tt_endjuv_to_init \u5f71\u54cd\u8425\u517b\u751f\u957f\u5230\u751f\u6b96\u751f\u957f\u8f6c\u6362\u3002"), u(r"\u751f\u7406\u89e3\u91ca\uff1a\u5f00\u82b1\u65f6\u5e8f\u548c\u704c\u6d46\u957f\u5ea6\u6539\u53d8\u751f\u7269\u91cf\u79ef\u7d2f\u7a97\u53e3\uff0c\u8fdb\u800c\u5f71\u54cd\u4ea7\u91cf\u3002")]
    add_bullets(slide, 7.05, 1.25, 5.35, 3.7, bullets, size=18)
    add_caption(slide, 0.9, 5.55, 5.7, u(r"\u56fe\uff1a\u7389\u7c73 grain_yield \u7684 Top ST \u53c2\u6570\u3002"))
    add_takeaway(slide, u(r"\u7389\u7c73\u4ea7\u91cf\u654f\u611f\u6027\u9ad8\u5ea6\u96c6\u4e2d\uff0c\u4e3b\u5bfc\u8def\u5f84\u662f\u7269\u5019/\u70ed\u65f6\u95f4\u8c03\u63a7\u3002"))
    record(u(r"\u7389\u7c73\u4ea7\u91cf\u654f\u611f\u6027\u7ed3\u679c"), bullets, u(r"\u7389\u7c73 grain_yield Top ST \u6761\u5f62\u56fe"), u(r"\u8bb2\u6e05\u695a\u6bcf\u4e2a\u53c2\u6570\u548c\u751f\u7406\u8fc7\u7a0b\u7684\u5173\u7cfb\uff0c\u907f\u514d\u53ea\u8bfb\u53c2\u6570\u540d\u3002"))

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, u(r"\u5c0f\u9ea6\u4ea7\u91cf\u654f\u611f\u6027\u7ed3\u679c"), u(r"grain_yield \u53d7\u7c7d\u7c92\u5927\u5c0f\u3001\u704c\u6d46\u548c\u7269\u5019\u5171\u540c\u5f71\u54cd"))
    add_picture(slide, figs["wheat_bar"], 0.8, 1.25, w=5.8)
    bullets = [u(r"max_grain_size \u8868\u5f81\u6f5c\u5728\u7c7d\u7c92\u5927\u5c0f\uff0c\u662f\u91cd\u8981 sink \u76f8\u5173\u53c2\u6570\u3002"), u(r"tt_start_grain_fill \u8fde\u63a5\u704c\u6d46\u542f\u52a8\u548c\u4ea7\u91cf\u5f62\u6210\u3002"), u(r"tt_end_of_juvenile\u3001tt_floral_initiation \u548c vern_sens \u5171\u540c\u5f71\u54cd\u53d1\u80b2\u7a97\u53e3\u3002"), u(r"\u5c0f\u9ea6 Top \u53c2\u6570 ST \u63a5\u8fd1\uff0c\u8bf4\u660e\u4e0d\u5b9c\u53ea\u5f3a\u8c03\u5355\u4e2a Top1\u3002")]
    add_bullets(slide, 7.05, 1.2, 5.35, 4.0, bullets, size=18)
    add_caption(slide, 0.9, 5.55, 5.7, u(r"\u56fe\uff1a\u5c0f\u9ea6 grain_yield \u7684 Top ST \u53c2\u6570\u3002"))
    add_takeaway(slide, u(r"\u5c0f\u9ea6\u4ea7\u91cf\u654f\u611f\u6027\u66f4\u5206\u6563\uff0c\u4f53\u73b0 sink potential\u3001grain filling \u4e0e phenology \u7684\u5171\u540c\u63a7\u5236\u3002"))
    record(u(r"\u5c0f\u9ea6\u4ea7\u91cf\u654f\u611f\u6027\u7ed3\u679c"), bullets, u(r"\u5c0f\u9ea6 grain_yield Top ST \u6761\u5f62\u56fe"), u(r"\u5f3a\u8c03\u5c0f\u9ea6\u7684\u654f\u611f\u6027\u7ed3\u6784\u6bd4\u7389\u7c73\u66f4\u590d\u6742\uff0c\u9002\u5408\u7528\u53c2\u6570\u7ec4\u6765\u89e3\u91ca\u3002"))

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, u(r"N64 vs N128 \u7a33\u5b9a\u6027\u9a8c\u8bc1"), u(r"ST \u6392\u540d\u7a33\u5b9a\uff0cN128 \u53ef\u4f5c\u4e3a\u4e3b\u7ed3\u679c"))
    add_picture(slide, figs["stability"], 0.75, 1.2, w=7.0)
    for i, (value, label) in enumerate([
        ("10 / 10", "Top3 overlap \u2265 2"),
        ("10 / 10", "Top5 overlap \u2265 4"),
        ("10 / 10", "ST Spearman \u2265 0.8"),
        ("0 / 10", "ST Spearman < 0.6"),
    ]):
        add_metric(slide, 8.15, 1.15 + i * 1.12, 3.7, 0.78, value, label, COLORS["green"] if i < 3 else COLORS["orange"])
    add_caption(slide, 0.85, 5.15, 7.0, u(r"\u56fe\uff1a\u5404 crop \u00d7 target \u7684 ST \u6392\u540d\u5728 N64 \u4e0e N128 \u95f4\u7684 Spearman \u76f8\u5173\u3002"))
    add_takeaway(slide, u(r"\u8bba\u6587\u4e2d\u5efa\u8bae\uff1aN128 \u4f5c\u4e3a\u4e3b\u7ed3\u679c\uff0cN64 \u4f5c\u4e3a\u7a33\u5b9a\u6027\u9a8c\u8bc1\uff1bS1 \u4ec5\u4f5c\u8f85\u52a9\u8bf4\u660e\u3002"))
    record(u(r"N64 vs N128 \u7a33\u5b9a\u6027\u9a8c\u8bc1"), [u(r"Top3 overlap \u22652\uff1a10/10"), u(r"Top5 overlap \u22654\uff1a10/10"), u(r"ST Spearman \u22650.8\uff1a10/10"), u(r"\u7ed3\u8bba\uff1aN128 \u4f5c\u4e3a\u4e3b\u7ed3\u679c")], u(r"ST Spearman \u7a33\u5b9a\u6027\u67f1\u72b6\u56fe + \u6307\u6807\u5361\u7247"), u(r"\u8fd9\u4e00\u9875\u7528\u4e8e\u652f\u6491\u201c\u7ed3\u679c\u53ef\u9760\u201d\uff1b\u8bf4\u660e S1 \u6709\u6ce2\u52a8\uff0c\u4e3b\u6587\u4e3b\u8981\u89e3\u91ca ST\u3002"))

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, u(r"\u7ed3\u8bba\u4e0e\u4e0b\u4e00\u6b65\u5de5\u4f5c"), u(r"\u5f62\u6210\u8bba\u6587\u4e3b\u7ebf\uff0c\u540c\u65f6\u4fdd\u7559\u8f93\u51fa\u53d8\u91cf\u8fb9\u754c"))
    left = [u(r"\u7ed3\u8bba 1\uff1a\u7389\u7c73 grain_yield \u4e3b\u8981\u7531\u7269\u5019/\u70ed\u65f6\u95f4\u53c2\u6570\u4e3b\u5bfc\u3002"), u(r"\u7ed3\u8bba 2\uff1a\u5c0f\u9ea6 grain_yield \u540c\u65f6\u53d7\u7c7d\u7c92\u5927\u5c0f\u3001\u704c\u6d46\u548c\u53d1\u80b2\u9636\u6bb5\u53c2\u6570\u5f71\u54cd\u3002"), u(r"\u7ed3\u8bba 3\uff1a\u5c0f\u9ea6\u548c\u7389\u7c73\u8868\u73b0\u51fa\u4e0d\u540c\u7684\u54c1\u79cd\u53c2\u6570\u654f\u611f\u6027\u7ed3\u6784\u3002"), u(r"\u7ed3\u8bba 4\uff1aN64/N128 \u5bf9\u6bd4\u8868\u660e ST \u6392\u540d\u7a33\u5b9a\uff0c\u7ed3\u679c\u5177\u5907\u8bba\u6587\u89e3\u91ca\u57fa\u7840\u3002")]
    right = [u(r"\u9650\u5236\uff1agrain_number\u3001grain_weight\u3001water_use_efficiency \u5f53\u524d\u672a\u8f93\u51fa\uff0c\u672a\u7eb3\u5165\u5206\u6790\u3002"), u(r"\u4e0b\u4e00\u6b65\uff1a\u5b8c\u5584 APSIM output/report\uff0c\u8865\u5145\u4ea7\u91cf\u6784\u6210\u548c\u6c34\u5206\u5229\u7528\u6548\u7387\u53d8\u91cf\u3002"), u(r"\u5efa\u8bae\uff1a\u5148\u7528 --limit 5 \u5c0f\u6837\u672c\u6d4b\u8bd5\u65b0\u589e\u53d8\u91cf\uff0c\u518d\u51b3\u5b9a\u662f\u5426\u91cd\u8dd1\u5b8c\u6574 Sobol\u3002")]
    add_bullets(slide, 0.75, 1.2, 5.95, 4.8, left, size=19)
    add_bullets(slide, 7.15, 1.2, 5.25, 4.3, right, size=19)
    add_takeaway(slide, u(r"\u603b\u7ed3\uff1a\u5f53\u524d\u7ed3\u679c\u652f\u6301\u201c\u4f5c\u7269\u5dee\u5f02\u5316\u7684 APSIM \u54c1\u79cd\u53c2\u6570\u654f\u611f\u6027\u7ed3\u6784\u201d\uff0c\u53ef\u8fdb\u5165\u8bba\u6587\u4e3b\u56fe\u548c\u7ed3\u679c\u5199\u4f5c\u3002"))
    record(u(r"\u7ed3\u8bba\u4e0e\u4e0b\u4e00\u6b65\u5de5\u4f5c"), left + right, u(r"\u53cc\u680f\u7ed3\u8bba\u4e0e\u8ba1\u5212"), u(r"\u6700\u540e\u6536\u675f\u5230\u8bba\u6587\u4e3b\u7ebf\uff1a\u73b0\u5728\u53ef\u5199\u4e3b\u7ed3\u679c\uff0c\u4f46\u7f3a\u5931\u53d8\u91cf\u4f5c\u4e3a\u540e\u7eed\u6269\u5c55\u3002"))

    # Remove the temporary Chinese test slide from the final deliverable so the
    # report remains within the requested 10-slide limit. The verification step
    # below still checks the final PPTX contains real Chinese characters.
    first_slide_rid = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(first_slide_rid)
    del prs.slides._sldIdLst[0]

    out_path = OUT_DIR / "APSIM_Classic_Sobol_CN_UTF8_verified.pptx"
    prs.save(out_path)
    return out_path, slide_info


def verify_pptx(pptx_path: Path) -> dict[str, object]:
    with zipfile.ZipFile(pptx_path, "r") as z:
        slide_xmls = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        xml_text = "\n".join(z.read(n).decode("utf-8", errors="strict") for n in slide_xmls)
        media_count = len([n for n in z.namelist() if n.startswith("ppt/media/")])
    chinese_chars = re.findall(r"[\u4e00-\u9fff]", xml_text)
    has_question_damage = "????" in xml_text or "??" + " APSIM" in xml_text
    has_test = ZH_TEST in xml_text
    has_font = FONT_FACE in xml_text
    return {
        "slide_xml_count": len(slide_xmls),
        "media_count": media_count,
        "chinese_char_count": len(chinese_chars),
        "has_question_damage": has_question_damage,
        "has_test": has_test,
        "has_font": has_font,
        "font_count": xml_text.count(FONT_FACE),
    }


def write_notes(slide_info: list[dict[str, object]]) -> Path:
    page_prefix = u(r"\u7b2c")
    page_suffix = u(r"\u9875\uff1a")
    bullet_heading = u(r"\u6bcf\u9875 bullet points")
    chart_heading = u(r"\u6bcf\u9875\u5efa\u8bae\u56fe\u8868")
    notes_heading = u(r"\u6bcf\u9875\u8bb2\u7a3f\u5907\u6ce8")
    final_heading = u(r"\u6700\u540e\u4e00\u9875\u603b\u7ed3\u53e5")
    lines = [u(r"# \u4e2d\u6587\u6c47\u62a5 PPT \u9010\u9875\u5185\u5bb9\u4e0e\u8bb2\u7a3f\u5907\u6ce8"), ""]
    for i, item in enumerate(slide_info, 1):
        lines.append(f"## {page_prefix} {i} {page_suffix}{item['title']}")
        lines.append("")
        lines.append(f"**{bullet_heading}**")
        for bullet in item["bullets"]:
            lines.append(f"- {bullet}")
        lines.append("")
        lines.append(f"**{chart_heading}**：{item['chart']}")
        lines.append("")
        lines.append(f"**{notes_heading}**：{item['notes']}")
        lines.append("")
    lines.append(f"## {final_heading}")
    lines.append("")
    lines.append(u(r"\u5f53\u524d\u7ed3\u679c\u652f\u6301\u201c\u4f5c\u7269\u5dee\u5f02\u5316\u7684 APSIM \u54c1\u79cd\u53c2\u6570\u654f\u611f\u6027\u7ed3\u6784\u201d\uff1a\u7389\u7c73\u4e3b\u8981\u7531\u7269\u5019/\u70ed\u65f6\u95f4\u53c2\u6570\u4e3b\u5bfc\uff0c\u5c0f\u9ea6\u5219\u7531\u7c7d\u7c92\u5927\u5c0f\u3001\u704c\u6d46\u548c\u7269\u5019\u54cd\u5e94\u5171\u540c\u63a7\u5236\uff1bN128 \u7a33\u5b9a\u6027\u9a8c\u8bc1\u8868\u660e ST \u6392\u540d\u53ef\u4f5c\u4e3a\u8bba\u6587\u4e3b\u7ed3\u679c\u3002"))
    path = OUT_DIR / "slide_titles_bullets_notes_cn.md"
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    figs = make_figures()
    pptx_path, slide_info = build_pptx(figs)
    verification = verify_pptx(pptx_path)
    if verification["has_question_damage"] or verification["chinese_char_count"] <= 0:
        raise RuntimeError(f"Chinese verification failed: {json.dumps(verification, ensure_ascii=False)}")
    notes_path = write_notes(slide_info)
    qa = {
        "pptx": str(pptx_path),
        "notes": str(notes_path),
        "font": FONT_FACE,
        "verification": verification,
        "encoding_policy": "UTF-8 source; native PowerPoint text boxes; East Asian font explicitly set.",
        "known_limitations": "Speaker notes are provided in Markdown because python-pptx does not stably author notes pages.",
    }
    qa_path = OUT_DIR / "qa_report.json"
    qa_path.write_text(json.dumps(qa, ensure_ascii=False, indent=2), encoding="utf-8")
    md_path = OUT_DIR / "qa_report.md"
    md_path.write_text(
        "\n".join(
            [
                "# QA report",
                "",
                f"- PPTX: `{pptx_path}`",
                f"- Notes: `{notes_path}`",
                f"- Font: `{FONT_FACE}`",
                f"- Slide XML count: {verification['slide_xml_count']}",
                f"- Embedded media count: {verification['media_count']}",
                f"- Chinese character count in slide XML: {verification['chinese_char_count']}",
                f"- Contains test Chinese sentence: {verification['has_test']}",
                f"- Contains question-mark damage: {verification['has_question_damage']}",
                f"- Microsoft YaHei occurrences: {verification['font_count']}",
                "",
                "All Chinese text is inserted as native PowerPoint text and each run has East Asian font set to Microsoft YaHei.",
            ]
        ),
        encoding="utf-8",
    )
    print(f"PPTX={pptx_path}")
    print(f"NOTES={notes_path}")
    print(f"QA={md_path}")
    print(json.dumps(verification, ensure_ascii=False))


if __name__ == "__main__":
    main()
